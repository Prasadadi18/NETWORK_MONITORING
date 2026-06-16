"""
generate_ppt_v2.py  —  Network Guardian ICWITE Presentation
Exact visual match to NewsMania Research style:
  • Dark navy background everywhere
  • Large cyan bold title at top-left
  • Thin cyan separator line under title
  • Neon-bordered dark cards
  • White body text, cyan bold keywords
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Colours (measured from NewsMania slides) ─────────────────────────────────
BG          = RGBColor(0x0A, 0x16, 0x28)   # slide background
BG_CARD     = RGBColor(0x0E, 0x20, 0x40)   # card fill (slightly lighter)
CYAN_TITLE  = RGBColor(0x38, 0xBD, 0xF8)   # large slide titles
CYAN_GLOW   = RGBColor(0x1E, 0x88, 0xE5)   # card border / separator line
CYAN_BOLD   = RGBColor(0x5B, 0xC8, 0xFF)   # bold keyword text inside cards
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCB, 0xD5, 0xE1)   # secondary body text
DARK_TEXT   = RGBColor(0x0A, 0x16, 0x28)
CYAN_LINE   = RGBColor(0x1E, 0x5F, 0xA8)   # thin separator line

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide):
    """Fill entire slide with dark navy."""
    sh = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG
    sh.line.fill.background()
    return sh


def rect(slide, l, t, w, h, fill=BG_CARD,
         border=None, bw=Pt(1.5), radius=False):
    """Draw a rectangle (optionally with rounded corners via XML)."""
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border:
        sh.line.color.rgb = border
        sh.line.width = bw
    else:
        sh.line.fill.background()

    if radius:
        # Patch in round corners via XML prstGeom
        sp = sh._element
        spPr = sp.find(qn('p:spPr'))
        old = spPr.find(qn('a:prstGeom'))
        if old is not None:
            spPr.remove(old)
        new_geom = etree.SubElement(spPr, qn('a:prstGeom'))
        new_geom.set('prst', 'roundRect')
        avLst = etree.SubElement(new_geom, qn('a:avLst'))
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj'); gd.set('fmla', 'val 30000')

    return sh


def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True, spacing=0):
    """Add a textbox."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    if spacing:
        p.space_before = Pt(spacing)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb


def multiline_txt(slide, lines, l, t, w, h,
                  size=13, color=WHITE, bold_color=CYAN_BOLD,
                  line_spacing=3):
    """
    lines: list of (text, bold_flag) tuples or plain strings.
    Strings starting with '• ' auto-get bullet styling.
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            line_txt, is_bold = item, False
        else:
            line_txt, is_bold = item

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(line_spacing)
        r = p.add_run(); r.text = line_txt
        r.font.size = Pt(size)
        r.font.bold = is_bold
        r.font.color.rgb = bold_color if is_bold else color
    return tb


def slide_header(slide, title, subtitle=None):
    """Large cyan title + thin separator line (NewsMania content slide pattern)."""
    txt(slide, title,
        Inches(0.4), Inches(0.15),
        Inches(12.5), Inches(0.9),
        size=38, bold=True, color=CYAN_TITLE)
    # thin separator line
    rect(slide,
         Inches(0.4), Inches(1.05),
         Inches(12.5), Inches(0.045),
         fill=CYAN_LINE, border=None)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.4), Inches(1.12),
            Inches(12.5), Inches(0.4),
            size=18, bold=False, color=WHITE)


def card(slide, l, t, w, h, title=None, title_icon=None,
         body_lines=None, title_size=20, body_size=13.5):
    """Dark card with neon cyan border — exact NewsMania card style."""
    rect(slide, l, t, w, h, fill=BG_CARD,
         border=CYAN_GLOW, bw=Pt(1.8), radius=True)
    y = t + Inches(0.18)
    if title_icon:
        txt(slide, title_icon, l + Inches(0.18), y,
            Inches(0.5), Inches(0.5), size=22, color=CYAN_BOLD)
        tx_l = l + Inches(0.68)
    else:
        tx_l = l + Inches(0.22)

    if title:
        txt(slide, title, tx_l, y,
            w - (tx_l - l) - Inches(0.15), Inches(0.5),
            size=title_size, bold=True, color=CYAN_BOLD)
        y += Inches(0.52)

    if body_lines:
        bx = slide.shapes.add_textbox(
            l + Inches(0.22), y,
            w - Inches(0.44), t + h - y - Inches(0.15))
        bx.text_frame.word_wrap = True
        first = True
        for item in body_lines:
            if isinstance(item, tuple):
                line_txt, is_bold = item
            else:
                line_txt, is_bold = item, False
            p = bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
            first = False
            p.space_before = Pt(4)
            r = p.add_run(); r.text = line_txt
            r.font.size = Pt(body_size)
            r.font.bold = is_bold
            r.font.color.rgb = CYAN_BOLD if is_bold else WHITE


def slide_num(slide, n, total):
    txt(slide, f"{n}/{total}",
        Inches(12.5), Inches(7.15), Inches(0.75), Inches(0.3),
        size=10, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


TOTAL = 14

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Circuit/glow decorative lines (simulated)
rect(s, 0, Inches(3.5), prs.slide_width, Inches(0.015), fill=CYAN_GLOW)
rect(s, 0, Inches(3.52), prs.slide_width, Inches(0.008),
     fill=RGBColor(0x05, 0x30, 0x60))

# Title
txt(s, "Network Guardian",
    Inches(0.5), Inches(0.6), Inches(12.3), Inches(1.3),
    size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "An Intelligent System That Automatically Protects Computer\nNetworks from Attacks Using Artificial Intelligence",
    Inches(0.5), Inches(1.95), Inches(12.3), Inches(0.9),
    size=20, bold=False, color=CYAN_TITLE, align=PP_ALIGN.CENTER)

# Divider
rect(s, Inches(2.5), Inches(3.0), Inches(8.33), Inches(0.04), fill=CYAN_GLOW)

# Authors
txt(s, "Arya Y P,  Adi Narayan Prasad G,  Ashrith S Jain,  Abhishek",
    Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
    size=18, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)

txt(s, "Under the guidance of  Prof. Nalina V",
    Inches(0.5), Inches(3.72), Inches(12.3), Inches(0.4),
    size=15, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Department of Information Science and Engineering\nB.M.S. College of Engineering, Bengaluru-560019, India",
    Inches(0.5), Inches(4.18), Inches(12.3), Inches(0.65),
    size=14, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Conference badge
rect(s, Inches(3.8), Inches(5.0), Inches(5.73), Inches(0.7),
     fill=RGBColor(0x0E, 0x20, 0x40), border=CYAN_GLOW, bw=Pt(1.5), radius=True)
txt(s, "ICWITE 2026  —  Accepted Paper",
    Inches(3.9), Inches(5.1), Inches(5.5), Inches(0.5),
    size=16, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)

# Key numbers
for i, (val, lbl) in enumerate([
        ("<1 s", "Attack\nDetection"),
        ("2.88 s", "Total Response\nTime"),
        ("99%", "Blocking\nSuccess Rate"),
]):
    lx = Inches(1.0) + i * Inches(3.78)
    rect(s, lx, Inches(5.9), Inches(3.33), Inches(1.35),
         fill=RGBColor(0x08, 0x18, 0x35), border=CYAN_GLOW, bw=Pt(1.5))
    txt(s, val, lx + Inches(0.1), Inches(5.95), Inches(3.13), Inches(0.6),
        size=30, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)
    txt(s, lbl, lx + Inches(0.1), Inches(6.55), Inches(3.13), Inches(0.6),
        size=12, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

slide_num(s, 1, TOTAL)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  THE CYBERSECURITY CHALLENGE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "The Cybersecurity Challenge")
slide_num(s, 2, TOTAL)

txt(s, "The Modern Network Threat Landscape",
    Inches(0.4), Inches(1.22), Inches(7.0), Inches(0.45),
    size=22, bold=False, color=WHITE)

bullets_left = [
    ("Cryptic Alerts:  ", True),
    ('Traditional IDS displays messages like "ICMP flood from 192.168.1.50 at 150 pps" with no context. Non-experts cannot act.', False),
    ("", False),
    ("Alert Fatigue:  ", True),
    ("Systems generate thousands of alerts daily. Security teams become desensitized and miss critical threats. [2]", False),
    ("", False),
    ("Manual Response Too Slow:  ", True),
    ("Attacks execute in milliseconds. Average breach detection: 207 days. Containment: +73 days. [1]", False),
    ("", False),
    ("Skills Gap:  ", True),
    ("Organizations lack trained personnel to interpret technical alerts and respond autonomously.", False),
]
bx = s.shapes.add_textbox(Inches(0.4), Inches(1.75), Inches(7.0), Inches(5.2))
bx.text_frame.word_wrap = True
first = True
i = 0
while i < len(bullets_left):
    item = bullets_left[i]
    if isinstance(item, tuple) and item[1] and i+1 < len(bullets_left):
        # bold label + normal continuation on same visual paragraph
        p = bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
        first = False
        p.space_before = Pt(5)
        r1 = p.add_run(); r1.text = "• " + item[0]
        r1.font.size = Pt(14); r1.font.bold = True
        r1.font.color.rgb = CYAN_BOLD
        r2 = p.add_run(); r2.text = bullets_left[i+1][0] if isinstance(bullets_left[i+1], tuple) else bullets_left[i+1]
        r2.font.size = Pt(14); r2.font.bold = False
        r2.font.color.rgb = WHITE
        i += 2
    elif isinstance(item, tuple) and not item[0]:
        if not first:
            p = bx.text_frame.add_paragraph(); p.space_before = Pt(2)
        i += 1
    else:
        i += 1

# Right card
card(s, Inches(7.7), Inches(1.22), Inches(5.3), Inches(5.8),
     title="Why Network Guardian?",
     body_lines=[
         ("Detects attacks in < 1 second", False),
         ("", False),
         ("Generates plain-English explanations\naccessible to non-experts", False),
         ("", False),
         ("Autonomously blocks attacker IPs\nvia OS routing commands", False),
         ("", False),
         ("Auto-reverts blocks after 60 s\n— false positives self-heal", False),
         ("", False),
         ("100% user comprehension vs.\n< 30% for raw technical alerts", False),
         ("", False),
         ("99.96% Wireshark validation\nagreement across all tests", False),
     ])


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  —  SOLUTION OVERVIEW  (3 cards — exactly like NewsMania slide 3)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
txt(s, "Network Guardian Solution Overview",
    Inches(0.4), Inches(0.2), Inches(12.5), Inches(1.0),
    size=42, bold=True, color=CYAN_TITLE, align=PP_ALIGN.CENTER)
rect(s, Inches(0.4), Inches(1.18), Inches(12.5), Inches(0.04), fill=CYAN_LINE)
slide_num(s, 3, TOTAL)

three = [
    ("🛡",  "Autonomous\nDetection",
     "Real-time monitoring of 5 attack types:\nICMP Flood, SYN Flood, UDP Flood,\nPort Scan, IP Fragmentation.\n\nSliding-window counters with\nsub-second latency (<1 s)."),
    ("🤖",  "AI-Powered\nExplanation",
     "Locally hosted TinyLlama LLM\n(via Ollama) generates plain-English\nexplanations for every detected attack.\n\nNo cloud — private, fast,\nand accessible to non-experts."),
    ("⚡",  "Autonomous\nMitigation",
     "Automatic OS-level route blocking\n(Windows: route add  |  Linux: ip route).\n\n60-second auto-revert bounds\nfalse-positive impact.\n99% blocking success rate."),
]
cw = Inches(3.9); ch = Inches(5.5)
for i, (icon, title, body) in enumerate(three):
    lx = Inches(0.5) + i * Inches(4.28)
    card(s, lx, Inches(1.35), cw, ch, title_icon=icon,
         title=title, title_size=22,
         body_lines=[(body, False)], body_size=14)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  —  PROPOSED SYSTEM ARCHITECTURE  (like NewsMania slide 4)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Proposed System Architecture")
slide_num(s, 4, TOTAL)

txt(s, "5-Layer Modular Design",
    Inches(0.4), Inches(1.2), Inches(6.5), Inches(0.45),
    size=22, bold=False, color=WHITE)

layers = [
    ("Layer 5", "Web Dashboard",      "Real-time Flask + Socket.IO browser UI."),
    ("Layer 4", "Mitigator",           "OS route blocking with 60 s auto-revert."),
    ("Layer 3", "LLM Analyzer",        "TinyLlama explains threats in plain English."),
    ("Layer 2", "Attack Detector",     "5 sliding-window detectors in parallel."),
    ("Layer 1", "Network Sniffer",     "Scapy packet capture, metadata extraction."),
]

bx = s.shapes.add_textbox(Inches(0.4), Inches(1.75), Inches(6.3), Inches(5.1))
bx.text_frame.word_wrap = True
first = True
for lbl, name, desc in layers:
    p = bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
    first = False
    p.space_before = Pt(6)
    r1 = p.add_run(); r1.text = f"• {name}  "
    r1.font.size = Pt(16); r1.font.bold = True; r1.font.color.rgb = CYAN_BOLD
    p2 = bx.text_frame.add_paragraph(); p2.space_before = Pt(1)
    r2 = p2.add_run(); r2.text = f"   {desc}"
    r2.font.size = Pt(13.5); r2.font.bold = False; r2.font.color.rgb = WHITE

# Right: stacked layer boxes
ly_colors = [
    RGBColor(0x0E,0x3A,0x6B),
    RGBColor(0x0B,0x4D,0x4A),
    RGBColor(0x2D,0x1B,0x6B),
    RGBColor(0x5A,0x1A,0x1A),
    RGBColor(0x0C,0x3B,0x1C),
]
box_h = Inches(0.82); gap = Inches(0.1)
sx = Inches(7.0); sy = Inches(1.35)
for i, ((lbl, name, desc), col) in enumerate(zip(layers, ly_colors)):
    top = sy + i*(box_h+gap)
    rect(s, sx, top, Inches(5.9), box_h, fill=col,
         border=CYAN_GLOW, bw=Pt(1.5))
    txt(s, f"Layer {5-i}  |  {name}",
        sx + Inches(0.2), top + Inches(0.15),
        Inches(5.5), Inches(0.55),
        size=15, bold=True, color=WHITE)

# Arrow at bottom showing data flow
rect(s, Inches(7.0), Inches(5.9), Inches(5.9), Inches(0.5),
     fill=RGBColor(0x05,0x10,0x25), border=CYAN_GLOW, bw=Pt(1))
txt(s, "Packet  →  Detect  →  Explain  →  Block  →  Dashboard",
    Inches(7.1), Inches(5.96), Inches(5.7), Inches(0.4),
    size=13, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)

# Shared state bus label
txt(s, "Shared Queue Bus: packet_queue  |  threat_queue  |  ui_log_queue",
    Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.4),
    size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  —  ATTACK DETECTION ENGINE  (left text + right table)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Attack Detection Engine")
slide_num(s, 5, TOTAL)

txt(s, "Five Independent Sliding-Window Detectors",
    Inches(0.4), Inches(1.22), Inches(12.5), Inches(0.45),
    size=20, bold=False, color=WHITE)

# Left: key points
left_pts = [
    ("SlidingWindowCounter", True),
    (":  Timestamps in a deque per source IP; evict entries older than the window; count remaining.", False),
    ("PortScanTracker", True),
    (":  Unique destination ports per IP in a 5 s window — distinguishes scan from flood.", False),
    ("Cooldown (60 s)", True),
    (":  Per (IP, AttackType) pair — prevents alert storms; different attack types fire independently.", False),
    ("Detection triggers", True),
    (":  4–10× above measured traffic baseline — near-zero false-positive rate (<0.01).", False),
]
bx = s.shapes.add_textbox(Inches(0.4), Inches(1.75), Inches(5.8), Inches(4.8))
bx.text_frame.word_wrap = True
first = True
i = 0
while i < len(left_pts):
    item = left_pts[i]
    if isinstance(item, tuple) and item[1] and i+1 < len(left_pts):
        p = bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
        first = False; p.space_before = Pt(8)
        r1 = p.add_run(); r1.text = "• " + item[0]
        r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = CYAN_BOLD
        r2 = p.add_run(); r2.text = left_pts[i+1][0]
        r2.font.size = Pt(14); r2.font.bold = False; r2.font.color.rgb = WHITE
        i += 2
    else:
        i += 1

# Right: detection parameters table
col_heads = ["Attack Type", "Threshold", "Window", "Cooldown"]
rows_data = [
    ["ICMP Flood",        "20 pps",    "1.0 s",  "60 s"],
    ["SYN Flood",         "20 pps",    "1.0 s",  "60 s"],
    ["UDP Flood",         "20 pps",    "1.0 s",  "60 s"],
    ["Port Scan",         "10 ports",  "5.0 s",  "60 s"],
    ["IP Fragmentation",  "10 pps",    "1.0 s",  "60 s"],
]
col_ws = [Inches(1.95), Inches(1.2), Inches(1.1), Inches(1.1)]
col_xs = [Inches(6.5)]
for cw in col_ws[:-1]:
    col_xs.append(col_xs[-1] + cw)
row_h = Inches(0.52)
ty = Inches(1.35)

for ci, (hdr, cw, cx) in enumerate(zip(col_heads, col_ws, col_xs)):
    rect(s, cx, ty, cw, row_h, fill=RGBColor(0x0D,0x34,0x5C),
         border=CYAN_GLOW, bw=Pt(1))
    txt(s, hdr, cx+Inches(0.05), ty+Inches(0.1), cw-Inches(0.1), row_h-Inches(0.1),
        size=13, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)

row_fills = [RGBColor(0x0E,0x20,0x40), RGBColor(0x08,0x18,0x30)]
for ri, row in enumerate(rows_data):
    ry = ty + (ri+1)*row_h
    for ci, (cell, cw, cx) in enumerate(zip(row, col_ws, col_xs)):
        rect(s, cx, ry, cw, row_h,
             fill=row_fills[ri%2], border=CYAN_GLOW, bw=Pt(0.8))
        c = CYAN_BOLD if ci == 0 else WHITE
        txt(s, cell, cx+Inches(0.05), ry+Inches(0.1),
            cw-Inches(0.1), row_h-Inches(0.1),
            size=13, bold=(ci==0), color=c, align=PP_ALIGN.CENTER)

# Detection latency results sub-table
txt(s, "Detection Latency — Measured Results (Table IV of paper)",
    Inches(6.5), Inches(4.6), Inches(6.5), Inches(0.4),
    size=13, bold=True, color=CYAN_TITLE)
lat_heads = ["Attack", "Min (s)", "Mean (s)", "Max (s)"]
lat_data  = [
    ["ICMP",  "0.11", "0.14", "0.19"],
    ["SYN",   "0.23", "0.27", "0.34"],
    ["UDP",   "0.12", "0.15", "0.21"],
    ["Scan",  "0.45", "0.51", "0.62"],
    ["Frag",  "0.18", "0.21", "0.27"],
]
lcw = [Inches(1.55), Inches(1.18), Inches(1.22), Inches(1.2)]
lcx = [Inches(6.5)]
for cw2 in lcw[:-1]:
    lcx.append(lcx[-1]+cw2)
lty = Inches(5.05)
for ci, (hdr, cw2, cx2) in enumerate(zip(lat_heads, lcw, lcx)):
    rect(s, cx2, lty, cw2, Inches(0.42), fill=RGBColor(0x05,0x25,0x50),
         border=CYAN_GLOW, bw=Pt(0.8))
    txt(s, hdr, cx2+Inches(0.04), lty+Inches(0.07), cw2-Inches(0.08), Inches(0.32),
        size=11.5, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)
for ri, row in enumerate(lat_data):
    ry2 = lty + (ri+1)*Inches(0.42)
    for ci, (cell, cw2, cx2) in enumerate(zip(row, lcw, lcx)):
        rect(s, cx2, ry2, cw2, Inches(0.42),
             fill=row_fills[ri%2], border=CYAN_GLOW, bw=Pt(0.6))
        txt(s, cell, cx2+Inches(0.04), ry2+Inches(0.08),
            cw2-Inches(0.08), Inches(0.3),
            size=11.5, bold=(ci==0), color=CYAN_BOLD if ci==0 else WHITE,
            align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  —  LLM ANALYZER  (2-panel like NewsMania lit review slides)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "LLM Analyzer: AI-Powered Explanation")
txt(s, "Local LLM Reasoning — TinyLlama via Ollama",
    Inches(0.4), Inches(1.12), Inches(12.5), Inches(0.38),
    size=18, bold=True, color=CYAN_TITLE, align=PP_ALIGN.CENTER)
slide_num(s, 6, TOTAL)

# Left card — Key Design
card(s, Inches(0.35), Inches(1.6), Inches(6.1), Inches(5.65),
     title_icon="🧠", title="Key Design",
     body_lines=[
         ("• TinyLlama 1.1B parameter LLM via Ollama", False),
         ("  running locally at localhost:11434", False),
         ("• Prompt includes:", False),
         ("  – Attacker IP, attack type, packet count,", False),
         ("  – Normal baseline (< 20 pps)", False),
         ("• Returns structured JSON:", False),
         ('  { "attack", "explanation", "mitigation" }', False),
         ("• temperature = 0.1  — near-deterministic", False),
         ("• 3-retry policy  |  30 s timeout per attempt", False),
         ("• Fallback reasoning if Ollama unreachable", False),
         ("• IP hallucination guard — corrects wrong IP", False),
         ("• 94% first-attempt JSON success rate", False),
         ("• Mean analysis time: 2.46 s  (85.4% of pipeline)", False),
     ], body_size=13)

# Right card — Prompt + Output example
card(s, Inches(6.7), Inches(1.6), Inches(6.3), Inches(5.65),
     title_icon="⚙", title="Prompt & Output Example", body_size=12)

# Code-style inner box
rect(s, Inches(6.95), Inches(2.42), Inches(5.8), Inches(2.05),
     fill=RGBColor(0x04,0x10,0x22), border=CYAN_GLOW, bw=Pt(1))
txt(s, "PROMPT SENT TO OLLAMA:",
    Inches(7.1), Inches(2.48), Inches(5.5), Inches(0.3),
    size=10, bold=True, color=CYAN_BOLD)
prompt_ex = (
    "Attack Evidence:\n"
    "  Source IP:    203.0.113.50\n"
    "  Attack Type:  SYN Flood\n"
    "  Count:        87 packets in 1 s\n"
    "  Baseline:     < 20 pps\n\n"
    "Respond with ONLY JSON:\n"
    '{ "attack": "...", "explanation": "...",\n'
    '  "mitigation": "route add ..." }'
)
txt(s, prompt_ex, Inches(7.1), Inches(2.78), Inches(5.5), Inches(1.65),
    size=10.5, color=RGBColor(0x80,0xFF,0x80))

rect(s, Inches(6.95), Inches(4.58), Inches(5.8), Inches(1.5),
     fill=RGBColor(0x04,0x10,0x22), border=CYAN_GLOW, bw=Pt(1))
txt(s, "LLM RESPONSE:",
    Inches(7.1), Inches(4.64), Inches(5.5), Inches(0.3),
    size=10, bold=True, color=CYAN_BOLD)
resp_ex = (
    '"attack": "SYN Flood",\n'
    '"explanation": "Attacker sends SYN packets\n'
    '  without ACK, exhausting TCP queue.\n'
    '  Blocking stops new connections.",\n'
    '"mitigation": "route add 203.0.113.50\n'
    '  mask 255.255.255.255 192.0.2.1"'
)
txt(s, resp_ex, Inches(7.1), Inches(4.95), Inches(5.5), Inches(1.1),
    size=10.5, color=RGBColor(0xFF,0xE0,0x80))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7  —  AUTONOMOUS MITIGATION  (left + right layout)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Autonomous Mitigation Layer")
slide_num(s, 7, TOTAL)

txt(s, "Security-First Design — No Shell Injection",
    Inches(0.4), Inches(1.18), Inches(6.5), Inches(0.45),
    size=20, bold=False, color=WHITE)

left_bullets = [
    ("Command Injection Prevention:", True,
     "LLM output NEVER passed to shell=True. IP validated\nwith Python ipaddress module before any OS call."),
    ("OS Block Command:", True,
     "Windows:  route add <IP> mask 255.255.255.255 192.0.2.1\nLinux:      ip route add blackhole <IP>/32"),
    ("Auto-Revert (60 s):", True,
     "Daemon thread reverts block — false positives self-heal.\nProtects against routing damage from incorrect decisions."),
    ("Duplicate Prevention:", True,
     "IP marked in active_blocks before route command.\nConcurrent threats for same IP serialised."),
    ("Safety Checks:", True,
     "Refuses to block loopback (dashboard stays accessible).\nRefuses to block multicast addresses."),
]
bx = s.shapes.add_textbox(Inches(0.4), Inches(1.73), Inches(6.2), Inches(5.3))
bx.text_frame.word_wrap = True
first = True
for title, _, body in left_bullets:
    p = bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
    first = False; p.space_before = Pt(7)
    r = p.add_run(); r.text = "• " + title
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = CYAN_BOLD
    p2 = bx.text_frame.add_paragraph(); p2.space_before = Pt(1)
    r2 = p2.add_run(); r2.text = "   " + body
    r2.font.size = Pt(13); r2.font.color.rgb = WHITE

# Right: status lifecycle card
card(s, Inches(6.9), Inches(1.3), Inches(6.1), Inches(2.5),
     title_icon="🔄", title="Mitigation Status Lifecycle",
     body_lines=[
         ("PENDING  →  APPLIED  →  auto-revert 60 s", False),
         ("PENDING  →  FAILED   →  removed (route error)", False),
         ("", False),
         ("198 / 200 successful blocks in testing", False),
         ("2 failures: 1 routing lock, 1 invalid LLM cmd", False),
         ("Auto-revert: 60.02 ± 0.09 s (100 tests)", False),
     ], body_size=13)

# Wireshark result box
card(s, Inches(6.9), Inches(3.95), Inches(6.1), Inches(3.1),
     title_icon="📡", title="Wireshark Validation",
     body_lines=[
         ("Before block:", True),
         ("  Continuous attack traffic with signatures.", False),
         ("After block:", True),
         ("  Zero packets from blocked IPs at target.", False),
         ("  100% traffic reduction confirmed.", False),
         ("After auto-revert:", True),
         ("  Traffic resumes — revert mechanism correct.", False),
         ("", False),
         ("99.96% Scapy ↔ Wireshark packet count agreement", False),
     ], body_size=12.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8  —  WEB DASHBOARD  (like NewsMania "Personalization: Hybrid Filtering")
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Real-Time Web Dashboard")
slide_num(s, 8, TOTAL)

txt(s, "Layer 5 — Flask + Socket.IO at http://localhost:5000",
    Inches(0.4), Inches(1.22), Inches(6.2), Inches(0.42),
    size=16, bold=False, color=WHITE)

dashboard_pts = [
    ("Packet Feed:", True,
     "Last 50 packets, colour-coded by protocol.\nScrolling live table updated every 500 ms."),
    ("Threat History:", True,
     "Last 10 ThreatEvents with full LLM reasoning\ndisplayed inline for each detected attack."),
    ("Active Blocks:", True,
     "Current blocked IPs with countdown timers\nshowing seconds until auto-revert."),
    ("Attack Injection:", True,
     "Browser buttons launch test attacks:\nICMP | SYN | UDP | Port Scan | Frag | ALL"),
    ("Controls:", True,
     "Flush All Blocks — removes all active routes.\nShutdown — graceful exit, flushes all blocks."),
]
bx = s.shapes.add_textbox(Inches(0.4), Inches(1.75), Inches(6.2), Inches(5.0))
bx.text_frame.word_wrap = True
first = True
for title, _, body in dashboard_pts:
    p = bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
    first = False; p.space_before = Pt(7)
    r = p.add_run(); r.text = "• " + title
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = CYAN_BOLD
    p2 = bx.text_frame.add_paragraph(); p2.space_before = Pt(1)
    r2 = p2.add_run(); r2.text = "   " + body
    r2.font.size = Pt(13); r2.font.color.rgb = WHITE

# Right side — event + API box
rect(s, Inches(6.9), Inches(1.35), Inches(6.1), Inches(2.5),
     fill=RGBColor(0x04,0x10,0x22), border=CYAN_GLOW, bw=Pt(1.5))
txt(s, "Socket.IO Events (real-time push)",
    Inches(7.05), Inches(1.42), Inches(5.8), Inches(0.35),
    size=13, bold=True, color=CYAN_BOLD)
events = (
    "packet_feed   — new packet batch (50 max)\n"
    "threat_feed   — updated threat history\n"
    "log_feed      — new system log lines\n"
    "status_update — active blocks count + IPs"
)
txt(s, events, Inches(7.05), Inches(1.82), Inches(5.8), Inches(1.85),
    size=12.5, color=WHITE)

rect(s, Inches(6.9), Inches(4.0), Inches(6.1), Inches(2.0),
     fill=RGBColor(0x04,0x10,0x22), border=CYAN_GLOW, bw=Pt(1.5))
txt(s, "REST API Endpoints",
    Inches(7.05), Inches(4.07), Inches(5.8), Inches(0.35),
    size=13, bold=True, color=CYAN_BOLD)
apis = (
    "GET  /api/status    — blocks, threats, queue size\n"
    "POST /api/attack    — launch test attack\n"
    "POST /api/flush     — remove all active blocks\n"
    "POST /api/shutdown  — graceful system exit"
)
txt(s, apis, Inches(7.05), Inches(4.45), Inches(5.8), Inches(1.45),
    size=12.5, color=WHITE)

card(s, Inches(6.9), Inches(6.1), Inches(6.1), Inches(1.15),
     body_lines=[
         ("Background emitter thread: reads shared_state every 500 ms → pushes to all browsers via WebSocket", False)],
     body_size=12.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9  —  TECHNICAL IMPLEMENTATION  (3 cards like NewsMania slide 10)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Technical Implementation")
slide_num(s, 9, TOTAL)

three_impl = [
    ("🐍", "Python Stack",
     "Python 3.10+\nScapy + Npcap (Windows)\n/ libpcap (Linux)\nFlask + Socket.IO\n\nTinyLlama-1.1B via Ollama\nOS routing commands\nWireshark for validation"),
    ("📐", "Module Architecture",
     "7 Python modules — 1,720 lines total\n\nmain.py         250 lines\nshared_state.py  80 lines\nnetwork_sniffer 120 lines\nattack_detector 350 lines\nllm_analyzer    280 lines\nmitigator       240 lines\nattack_injector 400 lines"),
    ("🔗", "Queue-Based Design",
     "Fault isolation:\n  Each layer independent\n\nPacket queue:   5,000 cap\nThreat queue:   100 cap\nUI log queue:   500 cap\n\n5 s buffer @ 1,000 pps\n\nSingle-responsibility principle\nthroughout all modules"),
]
cw2 = Inches(3.9); ch2 = Inches(5.3)
for i, (icon, title, body) in enumerate(three_impl):
    lx = Inches(0.5) + i * Inches(4.28)
    card(s, lx, Inches(1.5), cw2, ch2, title_icon=icon,
         title=title, title_size=20,
         body_lines=[(body, False)], body_size=13.5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10  —  EXPERIMENTAL RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Experimental Results")
slide_num(s, 10, TOTAL)

txt(s, "End-to-End Response Time & System Resource Usage",
    Inches(0.4), Inches(1.22), Inches(12.5), Inches(0.42),
    size=20, bold=False, color=WHITE)

# Left: E2E breakdown card
card(s, Inches(0.35), Inches(1.75), Inches(5.8), Inches(3.8),
     title_icon="⏱", title="Pipeline Timing Breakdown",
     body_lines=[
         ("Stage            Mean      Std      %", False),
         ("────────────────────────────────────", False),
         ("Detection        0.25 s    0.12 s   8.7%", False),
         ("LLM Analysis     2.46 s    0.35 s  85.4%", False),
         ("Mitigation       0.17 s    0.04 s   5.9%", False),
         ("────────────────────────────────────", False),
         ("Total            2.88 s    0.41 s  100%", False),
         ("", False),
         ("LLM dominates at 85.4%", True),
         ("Still 2.88 s vs. 5–30 min for human analysts", False),
     ], body_size=12.5)

# Right: system resources
card(s, Inches(6.4), Inches(1.75), Inches(6.6), Inches(3.8),
     title_icon="💻", title="System Resource Usage",
     body_lines=[
         ("CPU Usage:", True),
         ("  Idle:  2.1%  |  Under attack: 14.7%", False),
         ("  Peak (multi-vector): 28.3%", False),
         ("", False),
         ("Memory:", True),
         ("  Idle:  87 MB  |  Peak: 149 MB", False),
         ("", False),
         ("Packet Queue Depth:", True),
         ("  Idle:  0  |  Under attack: 450 avg", False),
         ("  Peak: 1,200  (well within 5,000 cap)", False),
         ("", False),
         ("Suitable for lightweight embedded appliances", False),
     ], body_size=13)

# Bottom row: key metrics
metrics = [
    ("<1 s", "Attack Detection\nLatency"),
    ("99%",  "Blocking\nSuccess Rate"),
    ("0.01", "False Positive\nRate (<)"),
    ("100%", "Traffic Reduction\nAfter Block"),
]
for i, (val, lbl) in enumerate(metrics):
    lx = Inches(0.35) + i * Inches(3.25)
    rect(s, lx, Inches(5.72), Inches(3.0), Inches(1.55),
         fill=RGBColor(0x08,0x1A,0x38), border=CYAN_GLOW, bw=Pt(1.5))
    txt(s, val, lx+Inches(0.1), Inches(5.78), Inches(2.8), Inches(0.72),
        size=32, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)
    txt(s, lbl, lx+Inches(0.1), Inches(6.52), Inches(2.8), Inches(0.65),
        size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11  —  WIRESHARK VALIDATION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Wireshark Validation Methodology")
slide_num(s, 11, TOTAL)

txt(s, "Ground-Truth Verification Across All Five Attack Types",
    Inches(0.4), Inches(1.22), Inches(12.5), Inches(0.42),
    size=20, bold=False, color=WHITE)

# Workflow card
card(s, Inches(0.35), Inches(1.72), Inches(5.8), Inches(5.45),
     title_icon="📋", title="10-Step Validation Workflow",
     body_lines=[
         ("1.  Launch Wireshark with capture filters", False),
         ("2.  Start Network Guardian", False),
         ("3.  Initiate attack from attack_injector.py", False),
         ("4.  Observe real-time capture", False),
         ("5.  Apply attack-specific Wireshark filters", False),
         ("6.  Confirm detection triggers in dashboard", False),
         ("7.  Validate LLM explanation accuracy", False),
         ("8.  Verify block via routing table", False),
         ("9.  Confirm traffic cessation in Wireshark", False),
         ("10. Wait 60 s — verify auto-revert correct", False),
     ], body_size=13.5)

# Agreement table
txt(s, "Packet Count Agreement (5 × 10-minute sessions)",
    Inches(6.4), Inches(1.72), Inches(6.5), Inches(0.38),
    size=13, bold=True, color=CYAN_TITLE)

sessions = [
    ["Session", "Scapy", "Wireshark", "Agreement"],
    ["1",  "12,487", "12,483", "99.97%"],
    ["2",  "18,923", "18,915", "99.96%"],
    ["3",  " 9,456", " 9,452", "99.96%"],
    ["4",  "15,234", "15,229", "99.97%"],
    ["5",  "20,891", "20,881", "99.95%"],
    ["Mean",     "—",      "—",   "99.96%"],
]
scw = [Inches(1.0), Inches(1.35), Inches(1.45), Inches(1.45)]
scx = [Inches(6.4)]
for cw3 in scw[:-1]:
    scx.append(scx[-1]+cw3)
sty = Inches(2.15)
for ri, row in enumerate(sessions):
    is_hdr = (ri == 0)
    is_mean = (ri == 6)
    for ci, (cell, cw3, cx3) in enumerate(zip(row, scw, scx)):
        ry3 = sty + ri * Inches(0.48)
        fill = (RGBColor(0x0D,0x34,0x5C) if is_hdr
                else RGBColor(0x0A,0x25,0x45) if is_mean
                else row_fills[ri%2])
        rect(s, cx3, ry3, cw3, Inches(0.47),
             fill=fill, border=CYAN_GLOW, bw=Pt(0.8))
        c = CYAN_BOLD if (is_hdr or is_mean or ci==3) else WHITE
        txt(s, cell, cx3+Inches(0.04), ry3+Inches(0.08),
            cw3-Inches(0.08), Inches(0.32),
            size=12, bold=is_hdr or is_mean, color=c, align=PP_ALIGN.CENTER)

# Wireshark filter box
rect(s, Inches(6.4), Inches(5.52), Inches(6.55), Inches(1.65),
     fill=RGBColor(0x04,0x10,0x22), border=CYAN_GLOW, bw=Pt(1.2))
txt(s, "Wireshark Filters Used",
    Inches(6.55), Inches(5.58), Inches(6.2), Inches(0.32),
    size=12, bold=True, color=CYAN_BOLD)
filters = (
    "ICMP Flood:   icmp.type==8\n"
    "SYN Flood:    tcp.flags.syn==1 && tcp.flags.ack==0\n"
    "Port Scan:    tcp.flags.syn==1  (unique dst ports)\n"
    "Fragmentation: ip.flags.mf==1 || ip.frag_offset>0"
)
txt(s, filters, Inches(6.55), Inches(5.95), Inches(6.2), Inches(1.18),
    size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12  —  COMPARISON WITH EXISTING SYSTEMS  (3-card like NewsMania slide 12)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Comparison with Existing Approaches")
slide_num(s, 12, TOTAL)

txt(s, "Network Guardian vs. Prior Intrusion Detection Systems  (Table VI of paper)",
    Inches(0.4), Inches(1.22), Inches(12.5), Inches(0.42),
    size=17, bold=False, color=WHITE)

# Comparison table
comp_heads = ["Approach", "Detection Basis", "Explainability", "Auto Mitigation", "Response Time"]
comp_data = [
    ["Snort [11]",           "Static rules",           "None",          "No",              "Manual (minutes)"],
    ["Anomaly IDS [12]",     "Behaviour baseline",     "None",          "No",              "Manual (minutes)"],
    ["DNN/ML IDS [13]",      "Deep learning",          "None (black box)","No",            "Manual (minutes)"],
    ["LIME / SHAP [14,15]",  "Post-hoc ML",            "Technical only", "No",             "Manual"],
    ["Network Guardian",     "Sliding window + LLM",   "Plain English", "Yes (60 s revert)","2.88 s"],
]
ccw = [Inches(2.1), Inches(2.15), Inches(2.15), Inches(2.0), Inches(2.05)]
ccx = [Inches(0.35)]
for cw4 in ccw[:-1]:
    ccx.append(ccx[-1]+cw4)
ch_h = Inches(0.5); cr_h = Inches(0.82)
cty = Inches(1.78)
for ci, (hdr, cw4, cx4) in enumerate(zip(comp_heads, ccw, ccx)):
    rect(s, cx4, cty, cw4, ch_h, fill=RGBColor(0x0D,0x34,0x5C),
         border=CYAN_GLOW, bw=Pt(1))
    txt(s, hdr, cx4+Inches(0.05), cty+Inches(0.08),
        cw4-Inches(0.1), ch_h-Inches(0.1),
        size=12, bold=True, color=CYAN_BOLD, align=PP_ALIGN.CENTER)

for ri, row in enumerate(comp_data):
    is_ours = (ri == 4)
    ry4 = cty + ch_h + ri*cr_h
    fill = RGBColor(0x05,0x28,0x18) if is_ours else row_fills[ri%2]
    bord = RGBColor(0x00,0xFF,0x80) if is_ours else CYAN_GLOW
    bw4 = Pt(2) if is_ours else Pt(0.8)
    for ci, (cell, cw4, cx4) in enumerate(zip(row, ccw, ccx)):
        rect(s, cx4, ry4, cw4, cr_h, fill=fill, border=bord, bw=bw4)
        c = RGBColor(0x80,0xFF,0x80) if (is_ours and ci>0) else (CYAN_BOLD if ci==0 else WHITE)
        txt(s, cell, cx4+Inches(0.06), ry4+Inches(0.1),
            cw4-Inches(0.12), cr_h-Inches(0.12),
            size=11.5, bold=is_ours, color=c, align=PP_ALIGN.CENTER)

txt(s, "Network Guardian (highlighted green) uniquely combines deterministic sub-second detection + plain-English LLM explanations + autonomous mitigation.",
    Inches(0.35), Inches(6.98), Inches(12.6), Inches(0.4),
    size=11.5, bold=True, color=RGBColor(0x80,0xFF,0x80))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13  —  CHALLENGES & LIMITATIONS  (3-card like NewsMania slide 12)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
slide_header(s, "Challenges & Limitations")
slide_num(s, 13, TOTAL)

three_ch = [
    ("⏱", "Technical",
     "LLM response randomness:\n~2% of responses had formatting issues.\nSolved with exponential backoff + retry.\n\nCPU-only inference: 2.46 s per explanation.\nGPU deployment would reduce to < 0.5 s.\n\nAdversarial disinformation evolves rapidly\n— threshold-based detection may need\ncontinuous recalibration."),
    ("⚖", "Ethical",
     "Risk of false positives permanently\nblocking legitimate traffic (mitigated\nby 60 s auto-revert).\n\nLLM explanations are AI-generated —\nshould not replace human judgment\nfor critical decisions.\n\nBalancing autonomous action with\nhuman oversight in production."),
    ("🌐", "Scope",
     "Currently tested on local loopback only\n(RFC 5737 TEST-NET source IPs).\n\nSingle-host deployment — no distributed\nor multi-network support yet.\n\nIPv6 not yet supported (Scapy filter\nand route commands are IPv4 only).\n\nFive attack types only — DNS amp,\nHTTP slowloris not yet covered."),
]
cw5 = Inches(3.9); ch5 = Inches(5.3)
for i, (icon, title, body) in enumerate(three_ch):
    lx = Inches(0.5) + i * Inches(4.28)
    card(s, lx, Inches(1.5), cw5, ch5, title_icon=icon,
         title=title, title_size=20,
         body_lines=[(body, False)], body_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14  —  CONCLUSION & FUTURE SCOPE  (2-panel like NewsMania slide 13)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
txt(s, "Conclusion & Future Scope",
    Inches(0.4), Inches(0.18), Inches(12.5), Inches(0.9),
    size=42, bold=True, color=CYAN_TITLE, align=PP_ALIGN.CENTER)
rect(s, Inches(0.4), Inches(1.08), Inches(12.5), Inches(0.04), fill=CYAN_LINE)
slide_num(s, 14, TOTAL)

# Left — Summary
card(s, Inches(0.35), Inches(1.25), Inches(6.1), Inches(6.0),
     title_icon="🛡", title="Summary",
     body_lines=[
         ("Network Guardian demonstrates that integrating Local LLMs into network security is feasible and valuable.", False),
         ("", False),
         ("✓  Detects attacks in < 1 second", False),
         ("✓  Explains them in 2.5 seconds", False),
         ("✓  Blocks in 0.17 seconds with 60 s auto-revert", False),
         ("✓  99% blocking success rate", False),
         ("✓  < 0.01 false positive rate over 100 hours", False),
         ("✓  99.96% Wireshark validation agreement", False),
         ("✓  100% user comprehension of AI explanations\n    vs. < 30% for raw technical alerts", False),
         ("", False),
         ("Makes advanced security accessible to non-experts.\nValuable educational tool for learning attack\npatterns through AI-generated explanations.", False),
     ], body_size=13.5)

# Right — Future Roadmap
card(s, Inches(6.7), Inches(1.25), Inches(6.3), Inches(6.0),
     title_icon="🚀", title="Future Roadmap",
     body_lines=[
         ("GPU-Accelerated LLM:", True),
         ("  Llama-3 / Mistral on GPU → < 0.5 s analysis.", False),
         ("", False),
         ("Adaptive Thresholds:", True),
         ("  EWMA + Bayesian methods for dynamic\n  traffic baseline recalibration.", False),
         ("", False),
         ("Multi-Host Deployment:", True),
         ("  Distributed agents with centralised control.", False),
         ("", False),
         ("SIEM Integration:", True),
         ("  Export ThreatEvents to Elasticsearch / Splunk.", False),
         ("", False),
         ("Threat Intelligence Feeds:", True),
         ("  VirusTotal + MITRE ATT&CK + PostgreSQL.", False),
         ("", False),
         ("Extended Attack Coverage:", True),
         ("  DNS amplification, HTTP slowloris,\n  multiple LLM workers for parallelism.", False),
     ], body_size=13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15  —  QUESTIONS  (like NewsMania Q&A slide)  — NOT COUNTED IN 14
# But let's add it as bonus slide 15 (the user has 14 content slides + title)
# Actually let me include it — the NewsMania had 14 slides total including Q&A
# Let me rethink: slides 1-14 already have content; add Q&A as 15, update total
# Actually keep TOTAL=14 and skip slide_num on this one since it's the Q&A closer
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)

# Decorative corner dots
for lx, ty in [(Inches(0.3), Inches(0.2)), (Inches(12.6), Inches(0.2)),
                (Inches(0.3), Inches(7.1)), (Inches(12.6), Inches(7.1))]:
    rect(s, lx, ty, Inches(0.12), Inches(0.12), fill=CYAN_GLOW)

# Neon border frame
rect(s, Inches(0.35), Inches(0.3), Inches(12.63), Inches(6.9),
     fill=BG, border=CYAN_GLOW, bw=Pt(2))

# Polygon question mark (simulated with text)
txt(s, "?",
    Inches(4.5), Inches(0.8), Inches(4.33), Inches(2.8),
    size=140, bold=True, color=CYAN_TITLE, align=PP_ALIGN.CENTER)

txt(s, "Questions?",
    Inches(0.5), Inches(3.75), Inches(12.3), Inches(1.0),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Thank you for your attention.",
    Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
    size=24, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

rect(s, Inches(2.5), Inches(5.55), Inches(8.33), Inches(0.04), fill=CYAN_LINE)

txt(s, "Arya Y P  |  Adi Narayan Prasad G  |  Ashrith S Jain  |  Abhishek",
    Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.4),
    size=16, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Department of ISE  |  B.M.S. College of Engineering, Bengaluru",
    Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.38),
    size=14, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
out = r"c:\Users\adina\OneDrive\Desktop\cnfinal\NetworkGuardian_ICWITE_v2.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
