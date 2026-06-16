"""
generate_ppt_final.py
Network Guardian ICWITE Presentation — NewsMania style + 3D cyber images.
15 content slides + Q&A = 16 slides total.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import os

IMG = r"c:\Users\adina\OneDrive\Desktop\cnfinal\ppt_images"
OUT = r"c:\Users\adina\OneDrive\Desktop\cnfinal\NetworkGuardian_ICWITE_Final.pptx"

# ── Colours ───────────────────────────────────────────────────────────────────
BG         = RGBColor(0x0A,0x16,0x28)
BG_CARD    = RGBColor(0x0E,0x20,0x40)
CYAN_TITLE = RGBColor(0x38,0xBD,0xF8)
CYAN_GLOW  = RGBColor(0x1E,0x88,0xE5)
CYAN_BOLD  = RGBColor(0x5B,0xC8,0xFF)
WHITE      = RGBColor(0xFF,0xFF,0xFF)
LIGHT      = RGBColor(0xCB,0xD5,0xE1)
CYAN_LINE  = RGBColor(0x1E,0x5F,0xA8)
GREEN      = RGBColor(0x00,0xFF,0x80)
RED        = RGBColor(0xFF,0x44,0x44)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
BLANK = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────────────────────
def bg(slide):
    sh = slide.shapes.add_shape(1,0,0,prs.slide_width,prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = BG
    sh.line.fill.background(); return sh

def rect(slide,l,t,w,h,fill=BG_CARD,border=None,bw=Pt(1.5),radius=False):
    sh = slide.shapes.add_shape(1,l,t,w,h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border: sh.line.color.rgb=border; sh.line.width=bw
    else: sh.line.fill.background()
    if radius:
        sp=sh._element; spPr=sp.find(qn('p:spPr'))
        old=spPr.find(qn('a:prstGeom'))
        if old is not None: spPr.remove(old)
        g=etree.SubElement(spPr,qn('a:prstGeom')); g.set('prst','roundRect')
        av=etree.SubElement(g,qn('a:avLst')); gd=etree.SubElement(av,qn('a:gd'))
        gd.set('name','adj'); gd.set('fmla','val 25000')
    return sh

def txt(slide,text,l,t,w,h,size=14,bold=False,color=WHITE,
        align=PP_ALIGN.LEFT,italic=False,wrap=True):
    tb=slide.shapes.add_textbox(l,t,w,h)
    tf=tb.text_frame; tf.word_wrap=wrap
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text
    r.font.size=Pt(size); r.font.bold=bold
    r.font.italic=italic; r.font.color.rgb=color
    return tb

def img(slide,fname,l,t,w,h):
    path=os.path.join(IMG,fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path,l,t,w,h)

def header(slide,title,subtitle=None):
    txt(slide,title,Inches(0.4),Inches(0.12),Inches(12.5),Inches(0.88),
        size=38,bold=True,color=CYAN_TITLE)
    rect(slide,Inches(0.4),Inches(1.02),Inches(12.5),Inches(0.04),
         fill=CYAN_LINE,border=None)
    if subtitle:
        txt(slide,subtitle,Inches(0.4),Inches(1.1),Inches(12.5),Inches(0.38),
            size=17,color=WHITE)

def snum(slide,n,tot):
    txt(slide,f"{n} / {tot}",Inches(12.4),Inches(7.15),Inches(0.8),Inches(0.28),
        size=10,color=LIGHT,align=PP_ALIGN.RIGHT)

def card(slide,l,t,w,h,title=None,icon=None,lines=None,ts=19,bs=13):
    rect(slide,l,t,w,h,fill=BG_CARD,border=CYAN_GLOW,bw=Pt(1.8),radius=True)
    y=t+Inches(0.18)
    xl=l+Inches(0.22)
    if icon:
        txt(slide,icon,l+Inches(0.18),y,Inches(0.5),Inches(0.48),size=20,color=CYAN_BOLD)
        xl=l+Inches(0.68)
    if title:
        txt(slide,title,xl,y,w-(xl-l)-Inches(0.15),Inches(0.48),
            size=ts,bold=True,color=CYAN_BOLD)
        y+=Inches(0.52)
    if lines:
        bx=slide.shapes.add_textbox(l+Inches(0.22),y,
            w-Inches(0.44),t+h-y-Inches(0.12))
        bx.text_frame.word_wrap=True
        first=True
        for item in lines:
            lb,ib=(item,False) if isinstance(item,str) else item
            p=bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
            first=False; p.space_before=Pt(3.5)
            r=p.add_run(); r.text=lb
            r.font.size=Pt(bs); r.font.bold=ib
            r.font.color.rgb=CYAN_BOLD if ib else WHITE

def bullet_block(slide,items,l,t,w,h,size=13.5,heading_size=14):
    """items: list of (bold_label, body_text) tuples or plain strings."""
    bx=slide.shapes.add_textbox(l,t,w,h)
    bx.text_frame.word_wrap=True
    first=True
    for item in items:
        if isinstance(item,tuple) and len(item)==2:
            lbl,body=item
            p=bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
            first=False; p.space_before=Pt(6)
            r1=p.add_run(); r1.text="• "+lbl
            r1.font.size=Pt(heading_size); r1.font.bold=True
            r1.font.color.rgb=CYAN_BOLD
            if body:
                p2=bx.text_frame.add_paragraph(); p2.space_before=Pt(1)
                r2=p2.add_run(); r2.text="   "+body
                r2.font.size=Pt(size); r2.font.color.rgb=WHITE
        else:
            p=bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
            first=False; p.space_before=Pt(3)
            r=p.add_run(); r.text=str(item)
            r.font.size=Pt(size); r.font.color.rgb=WHITE

TOTAL=16

# ═══════════════════════════════════════════════════════════════════════════════
# S1  TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
rect(s,0,Inches(3.45),prs.slide_width,Inches(0.018),fill=CYAN_GLOW)

img(s,'img_shield.png',Inches(4.5),Inches(0.5),Inches(4.33),Inches(3.2))

txt(s,"Network Guardian",Inches(0.4),Inches(0.12),Inches(12.5),Inches(1.0),
    size=52,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"An Intelligent System That Automatically Protects Computer Networks\n"
    "from Attacks Using Artificial Intelligence",
    Inches(0.4),Inches(1.15),Inches(12.5),Inches(0.8),
    size=18,color=CYAN_TITLE,align=PP_ALIGN.CENTER)

txt(s,"Arya Y P  |  Adi Narayan Prasad G  |  Ashrith S Jain  |  Abhishek",
    Inches(0.4),Inches(3.75),Inches(12.5),Inches(0.48),
    size=17,bold=True,color=CYAN_BOLD,align=PP_ALIGN.CENTER)
txt(s,"Under the Guidance of  Prof. Nalina V",
    Inches(0.4),Inches(4.28),Inches(12.5),Inches(0.38),
    size=14,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"Department of Information Science & Engineering\n"
    "B.M.S. College of Engineering, Bengaluru – 560019",
    Inches(0.4),Inches(4.72),Inches(12.5),Inches(0.62),
    size=13,color=LIGHT,align=PP_ALIGN.CENTER)

rect(s,Inches(3.9),Inches(5.48),Inches(5.53),Inches(0.65),
     fill=RGBColor(0x0E,0x20,0x40),border=CYAN_GLOW,bw=Pt(1.5),radius=True)
txt(s,"ICWITE 2026 — Accepted Paper",Inches(3.95),Inches(5.58),
    Inches(5.43),Inches(0.45),size=15,bold=True,color=CYAN_BOLD,
    align=PP_ALIGN.CENTER)

for i,(val,lbl) in enumerate([("<1 s","Attack Detection"),
                               ("2.88 s","Total Response"),
                               ("99%","Blocking Success")]):
    lx=Inches(0.4)+i*Inches(4.18)
    rect(s,lx,Inches(6.25),Inches(3.98),Inches(1.1),
         fill=RGBColor(0x08,0x18,0x35),border=CYAN_GLOW,bw=Pt(1.5))
    txt(s,val,lx+Inches(0.1),Inches(6.28),Inches(3.78),Inches(0.58),
        size=28,bold=True,color=CYAN_BOLD,align=PP_ALIGN.CENTER)
    txt(s,lbl,lx+Inches(0.1),Inches(6.86),Inches(3.78),Inches(0.42),
        size=12,color=WHITE,align=PP_ALIGN.CENTER)
snum(s,1,TOTAL)

# ═══════════════════════════════════════════════════════════════════════════════
# S2  THE CYBERSECURITY CHALLENGE
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"The Cybersecurity Challenge")
snum(s,2,TOTAL)

txt(s,"The Modern Network Threat Landscape",
    Inches(0.4),Inches(1.18),Inches(7.0),Inches(0.42),size=21,color=WHITE)

bullet_block(s,[
    ("Cryptic Alerts",
     'IDS shows "ICMP flood from 192.168.1.50 at 150 pps" — no context.\n   Non-experts cannot act on raw technical messages.'),
    ("Alert Fatigue",
     "Systems generate thousands of alerts daily. Security teams\n   become desensitized and miss critical threats. [2]"),
    ("Manual Response Too Slow",
     "Attacks execute in milliseconds. Average breach detection: 207 days.\n   Containment adds 73 more days. [1]"),
    ("Cybersecurity Skills Gap",
     "Organizations lack trained personnel to interpret alerts\n   and respond appropriately under pressure."),
],Inches(0.4),Inches(1.65),Inches(6.4),Inches(5.5))

rect(s,Inches(7.1),Inches(1.18),Inches(5.85),Inches(3.5),
     fill=BG_CARD,border=CYAN_GLOW,bw=Pt(1.8),radius=True)
txt(s,"Why We Need Network Guardian",
    Inches(7.25),Inches(1.28),Inches(5.55),Inches(0.42),
    size=16,bold=True,color=CYAN_BOLD)
wpoints=[
    "✓  Detects attacks in < 1 second",
    "✓  Generates plain-English AI explanations",
    "✓  Autonomously blocks attacker IPs via OS route",
    "✓  60-second auto-revert — false positives self-heal",
    "✓  99.96% Wireshark validation agreement",
    "✓  100% user comprehension vs <30% for raw alerts",
    "✓  Accessible to non-security-expert operators",
]
bx=s.shapes.add_textbox(Inches(7.25),Inches(1.78),Inches(5.5),Inches(2.8))
bx.text_frame.word_wrap=True
for i,p_ in enumerate(wpoints):
    p=bx.text_frame.paragraphs[0] if i==0 else bx.text_frame.add_paragraph()
    p.space_before=Pt(4)
    r=p.add_run(); r.text=p_; r.font.size=Pt(12.5)
    r.font.color.rgb=WHITE

# Key stat
rect(s,Inches(7.1),Inches(4.78),Inches(5.85),Inches(1.58),
     fill=RGBColor(0x08,0x18,0x35),border=CYAN_GLOW,bw=Pt(1.5),radius=True)
txt(s,'Research Question:',Inches(7.25),Inches(4.88),Inches(5.5),Inches(0.35),
    size=13,bold=True,color=CYAN_BOLD)
txt(s,'Can a local LLM be integrated into a live packet-capture pipeline\n'
    'to autonomously detect and mitigate network attacks with sub-second latency?',
    Inches(7.25),Inches(5.25),Inches(5.5),Inches(1.0),
    size=12,color=WHITE,italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# S3  RELATED WORK
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Related Work","Gap Analysis — Why existing systems are insufficient")
snum(s,3,TOTAL)

refs=[
    ("Snort [11] — Signature-Based",
     "Open-source rule-based detection. Cannot generalize to zero-day exploits\nor novel attack variants. No explanations provided."),
    ("Anomaly IDS [12] — Behaviour Baseline",
     "Establishes traffic baselines and flags deviations. High false-positive\nrate: 10–30%, making production deployment challenging."),
    ("DNN / ML IDS [13]",
     "Deep neural network classifiers achieve 95–99% accuracy on benchmarks.\nHowever, black-box — no explanations, limiting trust."),
    ("LIME / SHAP [14,15]",
     "Post-hoc ML explanation methods. Technical output only — not accessible\nto non-security-expert operators."),
    ("LLM for Security [16]",
     "Recent work explores LLMs for security tasks. Network Guardian\ndistinguishes itself: LLM for explanation, not detection."),
    ("AI-Based IDS [18–24]",
     "KNN, Naive Bayes, SVM — effective on benchmarks but no single solution\ncovers all attack types or provides autonomous mitigation."),
]
bw_r=Inches(5.85); bh_r=Inches(1.3)
for i,(title,body) in enumerate(refs):
    col_=i%2; row_=i//2
    lx=Inches(0.35)+col_*Inches(6.63)
    ty=Inches(1.58)+row_*(bh_r+Inches(0.12))
    rect(s,lx,ty,bw_r,bh_r,fill=BG_CARD,border=CYAN_GLOW,bw=Pt(1.5),radius=True)
    txt(s,title,lx+Inches(0.18),ty+Inches(0.1),bw_r-Inches(0.3),Inches(0.38),
        size=13,bold=True,color=CYAN_BOLD)
    txt(s,body,lx+Inches(0.18),ty+Inches(0.5),bw_r-Inches(0.3),bh_r-Inches(0.58),
        size=11.5,color=WHITE)

rect(s,Inches(0.35),Inches(5.88),Inches(12.63),Inches(0.6),
     fill=RGBColor(0x05,0x20,0x10),border=GREEN,bw=Pt(1.5),radius=True)
txt(s,"Gap Identified:  No existing system combines sub-second sliding-window detection + local LLM plain-English "
    "explanations + autonomous OS-level mitigation with auto-revert in a single lightweight Python pipeline.",
    Inches(0.5),Inches(5.95),Inches(12.3),Inches(0.5),
    size=12.5,bold=True,color=GREEN)

# ═══════════════════════════════════════════════════════════════════════════════
# S4  SOLUTION OVERVIEW  (3 neon cards)
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
txt(s,"Network Guardian Solution Overview",
    Inches(0.4),Inches(0.15),Inches(12.5),Inches(0.95),
    size=40,bold=True,color=CYAN_TITLE,align=PP_ALIGN.CENTER)
rect(s,Inches(0.4),Inches(1.12),Inches(12.5),Inches(0.04),fill=CYAN_LINE)
snum(s,4,TOTAL)

three=[
    ("🛡","Autonomous Detection",
     "Real-time monitoring of 5 attack types:\n"
     "  • ICMP Flood   • SYN Flood   • UDP Flood\n"
     "  • Port Scan    • IP Fragmentation\n\n"
     "Sliding-window counters per source IP.\n"
     "Sub-second detection latency (<1 s).\n"
     "Cooldown 60 s per (IP, attack type) pair.\n"
     "< 0.01 false-positive rate over 100 hours."),
    ("🤖","AI-Powered Explanation",
     "Locally hosted TinyLlama 1.1B via Ollama.\n"
     "No cloud — private, offline, fast.\n\n"
     "Structured JSON prompt → plain-English output:\n"
     '  { "attack", "explanation", "mitigation" }\n\n'
     "temperature=0.1 — near-deterministic.\n"
     "94% first-attempt success rate.\n"
     "Mean analysis time: 2.46 s."),
    ("⚡","Autonomous Mitigation",
     "OS-level route blocking:\n"
     "  Windows: route add <IP> mask ... 192.0.2.1\n"
     "  Linux:   ip route add blackhole <IP>/32\n\n"
     "No shell=True — injection-proof.\n"
     "IP validated before every command.\n"
     "60-second auto-revert daemon thread.\n"
     "99% blocking success (198/200 tests)."),
]
for i,(icon,title,body) in enumerate(three):
    lx=Inches(0.42)+i*Inches(4.3)
    card(s,lx,Inches(1.28),Inches(4.1),Inches(6.0),
         title=title,icon=icon,lines=[(body,False)],ts=20,bs=13)

# ═══════════════════════════════════════════════════════════════════════════════
# S5  SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Proposed System Architecture","5-Layer Modular Queue-Based Pipeline")
snum(s,5,TOTAL)

txt(s,"5-Layer Modular Design",
    Inches(0.4),Inches(1.2),Inches(5.8),Inches(0.42),size=21,color=WHITE)

layers=[
    ("Layer 5","Web Dashboard",  "Flask+SocketIO real-time browser UI","#0E3A6B"),
    ("Layer 4","Mitigator",      "OS route block, 60 s auto-revert",   "#0B4D4A"),
    ("Layer 3","LLM Analyzer",   "TinyLlama explains threats in plain English","#2D1B6B"),
    ("Layer 2","Attack Detector","5 sliding-window detectors in parallel","#5A1A1A"),
    ("Layer 1","Network Sniffer","Scapy packet capture, metadata extraction","#0C3B1C"),
]
bx=s.shapes.add_textbox(Inches(0.4),Inches(1.72),Inches(5.8),Inches(5.0))
bx.text_frame.word_wrap=True
first=True
for _,name,desc,_ in layers:
    p=bx.text_frame.paragraphs[0] if first else bx.text_frame.add_paragraph()
    first=False; p.space_before=Pt(7)
    r1=p.add_run(); r1.text="• "+name+"  "
    r1.font.size=Pt(16); r1.font.bold=True; r1.font.color.rgb=CYAN_BOLD
    p2=bx.text_frame.add_paragraph(); p2.space_before=Pt(1)
    r2=p2.add_run(); r2.text="   "+desc
    r2.font.size=Pt(13); r2.font.color.rgb=WHITE

# 3D stack image
img(s,'img_stack.png',Inches(6.5),Inches(1.18),Inches(6.5),Inches(4.6))

# Shared queue bus
rect(s,Inches(0.35),Inches(6.6),Inches(12.63),Inches(0.68),
     fill=RGBColor(0x05,0x18,0x30),border=CYAN_GLOW,bw=Pt(1.2))
txt(s,"Shared Queue Bus:  packet_queue (5000)  |  threat_queue (100)  |  "
    "ui_log_queue (500)  |  active_blocks dict  |  threat_history list",
    Inches(0.5),Inches(6.68),Inches(12.3),Inches(0.52),
    size=12,bold=True,color=CYAN_BOLD,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# S6  ATTACK DETECTION ENGINE
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Attack Detection Engine",
       "Five independent sliding-window detectors — sub-second latency")
snum(s,6,TOTAL)

# Radar image
img(s,'img_radar.png',Inches(0.3),Inches(1.3),Inches(5.0),Inches(4.5))

# Detection params table
txt(s,"Detection Algorithm Parameters  (Table I of paper)",
    Inches(5.5),Inches(1.28),Inches(7.5),Inches(0.38),
    size=14,bold=True,color=CYAN_TITLE)

col_h=["Attack Type","Threshold","Window","Cooldown"]
rows_=[["ICMP Flood","20 pps","1.0 s","60 s"],
       ["SYN Flood","20 pps","1.0 s","60 s"],
       ["UDP Flood","20 pps","1.0 s","60 s"],
       ["Port Scan","10 ports","5.0 s","60 s"],
       ["IP Fragmentation","10 pps","1.0 s","60 s"]]
cws=[Inches(2.1),Inches(1.2),Inches(1.1),Inches(1.1)]
cxs=[Inches(5.5)]
for cw_ in cws[:-1]: cxs.append(cxs[-1]+cw_)
rh=Inches(0.5); ty_=Inches(1.72)
for ci,(h_,cw_,cx_) in enumerate(zip(col_h,cws,cxs)):
    rect(s,cx_,ty_,cw_,rh,fill=RGBColor(0x0D,0x34,0x5C),border=CYAN_GLOW,bw=Pt(1))
    txt(s,h_,cx_+Inches(0.05),ty_+Inches(0.1),cw_-Inches(0.1),rh-Inches(0.1),
        size=12.5,bold=True,color=CYAN_BOLD,align=PP_ALIGN.CENTER)
rfc=[RGBColor(0x0E,0x20,0x40),RGBColor(0x08,0x18,0x30)]
for ri,row in enumerate(rows_):
    ry_=ty_+(ri+1)*rh
    for ci,(cell,cw_,cx_) in enumerate(zip(row,cws,cxs)):
        rect(s,cx_,ry_,cw_,rh,fill=rfc[ri%2],border=CYAN_GLOW,bw=Pt(0.8))
        txt(s,cell,cx_+Inches(0.05),ry_+Inches(0.1),
            cw_-Inches(0.1),rh-Inches(0.1),
            size=12,bold=(ci==0),color=CYAN_BOLD if ci==0 else WHITE,
            align=PP_ALIGN.CENTER)

# Detection latency sub-table
txt(s,"Detection Latency Results  (Table IV)",
    Inches(5.5),Inches(4.42),Inches(7.5),Inches(0.35),
    size=13,bold=True,color=CYAN_TITLE)
lh=["Attack","Min (s)","Mean (s)","Max (s)"]
ld=[["ICMP","0.11","0.14","0.19"],["SYN","0.23","0.27","0.34"],
    ["UDP","0.12","0.15","0.21"],["Scan","0.45","0.51","0.62"],
    ["Frag","0.18","0.21","0.27"]]
lw=[Inches(1.5),Inches(1.18),Inches(1.22),Inches(1.15)]
lx2=[Inches(5.5)]
for lw_ in lw[:-1]: lx2.append(lx2[-1]+lw_)
lty=Inches(4.82)
lrh=Inches(0.45)
for ci,(h_,cw_,cx_) in enumerate(zip(lh,lw,lx2)):
    rect(s,cx_,lty,cw_,lrh,fill=RGBColor(0x05,0x25,0x50),border=CYAN_GLOW,bw=Pt(0.8))
    txt(s,h_,cx_+Inches(0.04),lty+Inches(0.08),cw_-Inches(0.08),lrh-Inches(0.1),
        size=11,bold=True,color=CYAN_BOLD,align=PP_ALIGN.CENTER)
for ri,row in enumerate(ld):
    ry_=lty+(ri+1)*lrh
    for ci,(cell,cw_,cx_) in enumerate(zip(row,lw,lx2)):
        rect(s,cx_,ry_,cw_,lrh,fill=rfc[ri%2],border=CYAN_GLOW,bw=Pt(0.6))
        txt(s,cell,cx_+Inches(0.04),ry_+Inches(0.08),
            cw_-Inches(0.08),lrh-Inches(0.1),
            size=11,bold=(ci==0),color=CYAN_BOLD if ci==0 else WHITE,
            align=PP_ALIGN.CENTER)

# 3D attack chart
img(s,'img_attacks3d.png',Inches(5.5),Inches(5.82),Inches(7.5),Inches(1.55))

# ═══════════════════════════════════════════════════════════════════════════════
# S7  LLM ANALYZER
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"LLM Analyzer  —  AI-Powered Explanation",
       "TinyLlama 1.1B via Ollama · Local · Private · No cloud")
snum(s,7,TOTAL)

# Neural net image
img(s,'img_neural.png',Inches(0.3),Inches(1.35),Inches(5.5),Inches(3.8))

# Right: key design
card(s,Inches(6.05),Inches(1.3),Inches(7.0),Inches(2.6),
     title="Key Design Features",icon="🧠",
     lines=[
         ("• TinyLlama-1.1B via Ollama at localhost:11434",False),
         ("• temperature=0.1  |  num_predict=200 (bounded)",False),
         ("• 3-retry policy with 30 s timeout per attempt",False),
         ("• Fallback reasoning if Ollama unreachable",False),
         ("• IP hallucination guard — auto-corrects wrong IP",False),
         ("• 94% first-attempt JSON success  |  Mean: 2.46 s",False),
     ],ts=16,bs=12.5)

# Prompt example box
rect(s,Inches(6.05),Inches(4.02),Inches(7.0),Inches(3.25),
     fill=RGBColor(0x04,0x10,0x22),border=CYAN_GLOW,bw=Pt(1.5))
txt(s,"Prompt → Ollama → JSON Response",
    Inches(6.2),Inches(4.1),Inches(6.7),Inches(0.35),
    size=13,bold=True,color=CYAN_BOLD)
txt(s,'INPUT:\n'
    '  Source IP:  203.0.113.50\n'
    '  Attack:     SYN Flood\n'
    '  Count:      87 pps  (baseline < 20 pps)\n\n'
    'OUTPUT (JSON):\n'
    '  "attack":      "SYN Flood"\n'
    '  "explanation": "Attacker sends SYN packets without\n'
    '                  completing handshake, exhausting TCP\n'
    '                  connection queue of 128-512 slots."\n'
    '  "mitigation":  "route add 203.0.113.50 mask\n'
    '                  255.255.255.255 192.0.2.1"',
    Inches(6.2),Inches(4.5),Inches(6.7),Inches(2.7),
    size=11,color=RGBColor(0x80,0xFF,0x80))

# Bottom: network flow image
img(s,'img_network.png',Inches(0.3),Inches(5.28),Inches(5.5),Inches(2.0))

# ═══════════════════════════════════════════════════════════════════════════════
# S8  AUTONOMOUS MITIGATION
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Autonomous Mitigation Layer",
       "Security-first design — injection-proof OS-level route blocking")
snum(s,8,TOTAL)

bullet_block(s,[
    ("Command Injection Prevention",
     "LLM output NEVER passed to shell=True.\n   IP validated with Python ipaddress module before any OS call."),
    ("OS Block Command",
     "Windows:  route add <IP> mask 255.255.255.255 192.0.2.1\n"
     "   Linux:    ip route add blackhole <IP>/32\n"
     "   Gateway 192.0.2.1 = RFC 5737 TEST-NET (unreachable blackhole)."),
    ("Auto-Revert  (60 seconds)",
     "Daemon thread sleeps 60 s → runs route delete.\n   False positives self-heal — no permanent routing damage."),
    ("Safety Checks",
     "Refuses to block loopback (dashboard stays accessible).\n   Refuses to block multicast. subprocess.run(shell=False, timeout=10)."),
    ("Duplicate Block Prevention",
     "IP marked in active_blocks before route command.\n   Concurrent threats for same IP serialised under lock."),
],Inches(0.4),Inches(1.5),Inches(6.4),Inches(5.2))

card(s,Inches(7.0),Inches(1.3),Inches(6.0),Inches(2.25),
     title="Mitigation Status Lifecycle",icon="🔄",
     lines=[
         ("PENDING  →  APPLIED  →  auto-revert 60 s",False),
         ("PENDING  →  FAILED   →  removed (route error)",False),
         ("",False),
         ("198 / 200 successful blocks in testing",False),
         ("Auto-revert: 60.02 ± 0.09 s (100 tests)",False),
         ("100% traffic reduction verified by Wireshark",False),
     ],ts=16,bs=13)

card(s,Inches(7.0),Inches(3.65),Inches(6.0),Inches(3.1),
     title="Wireshark Block Verification",icon="📡",
     lines=[
         ("Before block:", True),
         ("  Continuous attack traffic with signatures visible.",False),
         ("After block:", True),
         ("  Zero packets from blocked IP — 100% reduction.",False),
         ("After auto-revert (60 s):", True),
         ("  Traffic resumes — revert mechanism confirmed correct.",False),
         ("",False),
         ("99.96% Scapy ↔ Wireshark packet count agreement",False),
         ("across five 10-minute validation sessions.",False),
     ],ts=16,bs=12.5)

# ═══════════════════════════════════════════════════════════════════════════════
# S9  WEB DASHBOARD
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Real-Time Web Dashboard",
       "Layer 5 — Flask + Socket.IO · http://localhost:5000 · Updates every 500 ms")
snum(s,9,TOTAL)

bullet_block(s,[
    ("Packet Feed",
     "Last 50 packets, colour-coded by protocol (ICMP/TCP/UDP).\n   Live scrolling table with timestamp."),
    ("Threat History",
     "Last 10 ThreatEvents with full LLM reasoning displayed inline."),
    ("Active Blocks + Timers",
     "Blocked IPs with countdown to 60-second auto-revert."),
    ("Attack Injection Buttons",
     "Browser controls: ICMP | SYN | UDP | Port Scan | Frag | ALL."),
    ("Control Actions",
     "Flush All Blocks — removes all active OS routes.\n   Shutdown — graceful exit, flushes blocks."),
],Inches(0.4),Inches(1.5),Inches(6.3),Inches(4.8))

rect(s,Inches(6.8),Inches(1.3),Inches(6.2),Inches(2.35),
     fill=RGBColor(0x04,0x10,0x22),border=CYAN_GLOW,bw=Pt(1.5))
txt(s,"Socket.IO Events (server → browser push)",
    Inches(6.95),Inches(1.38),Inches(5.9),Inches(0.35),
    size=13,bold=True,color=CYAN_BOLD)
txt(s,"packet_feed    — new packet batch (up to 50)\n"
    "threat_feed    — updated ThreatEvent list\n"
    "log_feed       — new system log messages\n"
    "status_update  — active_blocks count + IP list",
    Inches(6.95),Inches(1.78),Inches(5.9),Inches(1.8),
    size=12.5,color=WHITE)

rect(s,Inches(6.8),Inches(3.72),Inches(6.2),Inches(1.95),
     fill=RGBColor(0x04,0x10,0x22),border=CYAN_GLOW,bw=Pt(1.5))
txt(s,"REST API Endpoints",
    Inches(6.95),Inches(3.8),Inches(5.9),Inches(0.35),
    size=13,bold=True,color=CYAN_BOLD)
txt(s,"GET  /api/status    — blocks, threats, queue size\n"
    "POST /api/attack    — launch test attack (rate capped 1000 pps)\n"
    "POST /api/flush     — remove all active blocks\n"
    "POST /api/shutdown  — graceful exit",
    Inches(6.95),Inches(4.2),Inches(5.9),Inches(1.42),
    size=12.5,color=WHITE)

card(s,Inches(6.8),Inches(5.75),Inches(6.2),Inches(1.05),
     lines=[("Background emitter thread reads shared_state every 500 ms → "
             "pushes all 4 event types to connected browsers via WebSocket. "
             "Memory capped: 500 packets, 200 log lines, 10 threats.",False)],bs=12)

# ═══════════════════════════════════════════════════════════════════════════════
# S10  TECHNICAL IMPLEMENTATION  (3 cards)
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Technical Implementation")
snum(s,10,TOTAL)

three2=[
    ("🐍","Python Stack",
     "Language:      Python 3.10+\n"
     "Packet capture: Scapy + Npcap (Win)\n"
     "               libpcap (Linux)\n"
     "LLM inference:  TinyLlama-1.1B\n"
     "               via Ollama\n"
     "Web framework:  Flask + Socket.IO\n"
     "Mitigation:     OS routing commands\n"
     "Validation:     Wireshark"),
    ("📐","Module Architecture",
     "7 modules — 1,720 lines total\n\n"
     "main.py         250 lines\n"
     "shared_state.py  80 lines\n"
     "network_sniffer 120 lines\n"
     "attack_detector 350 lines\n"
     "llm_analyzer    280 lines\n"
     "mitigator       240 lines\n"
     "attack_injector 400 lines"),
    ("🔗","Queue-Based Design",
     "Fault isolation:\n  Each layer fully independent\n\n"
     "packet_queue   5,000 cap\n"
     "threat_queue     100 cap\n"
     "ui_log_queue     500 cap\n\n"
     "5-second buffer @ 1,000 pps\n\n"
     "Single-responsibility principle\n"
     "throughout all 7 modules"),
]
for i,(icon,title,body) in enumerate(three2):
    lx=Inches(0.42)+i*Inches(4.3)
    card(s,lx,Inches(1.28),Inches(4.1),Inches(6.0),
         title=title,icon=icon,lines=[(body,False)],ts=20,bs=13)

# ═══════════════════════════════════════════════════════════════════════════════
# S11  EXPERIMENTAL RESULTS
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Experimental Results",
       "Tested on Windows 11 | Npcap loopback | TinyLlama CPU inference")
snum(s,11,TOTAL)

img(s,'img_timing.png',Inches(0.3),Inches(1.3),Inches(8.8),Inches(4.2))

card(s,Inches(9.3),Inches(1.3),Inches(3.75),Inches(4.2),
     title="Key Numbers",icon="📊",
     lines=[
         ("Detection",True),("< 1 second",False),("",False),
         ("LLM Analysis",True),("2.46 s mean",False),("",False),
         ("Mitigation",True),("0.17 s mean",False),("",False),
         ("Total Pipeline",True),("2.88 s end-to-end",False),("",False),
         ("Blocking Success",True),("99%  (198/200)",False),("",False),
         ("False Positive Rate",True),("< 0.01",False),
     ],ts=15,bs=12.5)

# 4 metric badges
for i,(val,lbl) in enumerate([
        ("<1 s","Attack\nDetection"),("2.88 s","Total\nResponse"),
        ("99%","Blocking\nSuccess"),("0.01","False +ve\nRate (<)")]):
    lx_=Inches(0.3)+i*Inches(3.2)
    rect(s,lx_,Inches(5.62),Inches(3.05),Inches(1.6),
         fill=RGBColor(0x08,0x18,0x38),border=CYAN_GLOW,bw=Pt(1.5))
    txt(s,val,lx_+Inches(0.1),Inches(5.68),Inches(2.85),Inches(0.75),
        size=30,bold=True,color=CYAN_BOLD,align=PP_ALIGN.CENTER)
    txt(s,lbl,lx_+Inches(0.1),Inches(6.42),Inches(2.85),Inches(0.72),
        size=12,color=WHITE,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# S12  WIRESHARK VALIDATION
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Wireshark Validation",
       "Ground-truth verification — 99.96% Scapy ↔ Wireshark packet agreement")
snum(s,12,TOTAL)

img(s,'img_wireshark.png',Inches(0.3),Inches(1.28),Inches(7.5),Inches(3.5))

card(s,Inches(8.0),Inches(1.28),Inches(5.0),Inches(3.5),
     title="Session Agreement",icon="✅",
     lines=[
         ("Session 1: 12,487 vs 12,483 → 99.97%",False),
         ("Session 2: 18,923 vs 18,915 → 99.96%",False),
         ("Session 3:  9,456 vs  9,452 → 99.96%",False),
         ("Session 4: 15,234 vs 15,229 → 99.97%",False),
         ("Session 5: 20,891 vs 20,881 → 99.95%",False),
         ("",False),
         ("Mean Agreement: 99.96%",True),
     ],ts=16,bs=12)

# Workflow + filters
card(s,Inches(0.3),Inches(4.88),Inches(5.8),Inches(2.4),
     title="10-Step Workflow",icon="📋",
     lines=[
         ("1. Wireshark capture filters  →  2. Start Guardian",False),
         ("3. Inject attack  →  4. Observe live capture",False),
         ("5. Apply attack filter  →  6. Confirm detection",False),
         ("7. Validate LLM explanation accuracy",False),
         ("8. Verify block via routing table",False),
         ("9. Confirm traffic stops  →  10. Wait 60 s revert",False),
     ],ts=16,bs=12)

card(s,Inches(6.3),Inches(4.88),Inches(6.7),Inches(2.4),
     title="Wireshark Filters",icon="🔍",
     lines=[
         ("ICMP Flood:      icmp.type==8",False),
         ("SYN Flood:       tcp.flags.syn==1 && tcp.flags.ack==0",False),
         ("Port Scan:       tcp.flags.syn==1  (unique dst ports)",False),
         ("Fragmentation:   ip.flags.mf==1 || ip.frag_offset>0",False),
     ],ts=16,bs=12)

# ═══════════════════════════════════════════════════════════════════════════════
# S13  COMPARISON  (bar chart + table)
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Comparison with Existing Approaches",
       "Network Guardian vs. Snort / Anomaly IDS / DNN/ML / LIME+SHAP  (Table VI)")
snum(s,13,TOTAL)

img(s,'img_comparison.png',Inches(0.3),Inches(1.28),Inches(7.5),Inches(3.8))

# Right summary
card(s,Inches(8.0),Inches(1.28),Inches(5.0),Inches(3.8),
     title="What Makes Us Different",icon="★",
     lines=[
         ("Snort [11]",True),
         ("  ✗ No explainability  ✗ No auto-mitigation",False),
         ("Anomaly IDS [12]",True),
         ("  ✗ 10–30% false positives  ✗ No mitigation",False),
         ("DNN/ML IDS [13]",True),
         ("  ✗ Black-box  ✗ No explanation  ✗ No action",False),
         ("LIME/SHAP [14,15]",True),
         ("  ✗ Technical only  ✗ No mitigation",False),
         ("Network Guardian",True),
         ("  ✓ Plain English  ✓ Auto-block  ✓ 2.88 s  ✓ Revert",False),
     ],ts=16,bs=12)

# Bottom comparison mini-table
heads=["Approach","Explainability","Auto Mitigation","Response"]
data_=[["Snort","None","No","Manual"],
       ["Anomaly IDS","None","No","Manual"],
       ["DNN/ML IDS","None (black box)","No","Manual"],
       ["Network Guardian","Plain English","Yes (60s revert)","2.88 s"]]
cws2=[Inches(2.4),Inches(2.2),Inches(2.2),Inches(1.55)]
cxs2=[Inches(0.3)]
for cw_ in cws2[:-1]: cxs2.append(cxs2[-1]+cw_)
ty2=Inches(5.2); rh2=Inches(0.47)
for ci,(h_,cw_,cx_) in enumerate(zip(heads,cws2,cxs2)):
    rect(s,cx_,ty2,cw_,rh2,fill=RGBColor(0x0D,0x34,0x5C),border=CYAN_GLOW,bw=Pt(1))
    txt(s,h_,cx_+Inches(0.05),ty2+Inches(0.08),cw_-Inches(0.1),rh2-Inches(0.1),
        size=12,bold=True,color=CYAN_BOLD,align=PP_ALIGN.CENTER)
for ri,row in enumerate(data_):
    is_ours=(ri==3)
    ry_=ty2+(ri+1)*rh2
    fill_=RGBColor(0x05,0x28,0x18) if is_ours else rfc[ri%2]
    brd_=GREEN if is_ours else CYAN_GLOW
    for ci,(cell,cw_,cx_) in enumerate(zip(row,cws2,cxs2)):
        rect(s,cx_,ry_,cw_,rh2,fill=fill_,border=brd_,
             bw=Pt(2) if is_ours else Pt(0.8))
        txt(s,cell,cx_+Inches(0.05),ry_+Inches(0.08),
            cw_-Inches(0.1),rh2-Inches(0.1),
            size=11.5,bold=is_ours,
            color=GREEN if is_ours else WHITE,
            align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# S14  CHALLENGES & LIMITATIONS  (3 cards)
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
header(s,"Challenges & Limitations")
snum(s,14,TOTAL)

three3=[
    ("⏱","Technical Challenges",
     "LLM response randomness:\n"
     "  ~2% responses had formatting errors.\n"
     "  Solved: exponential backoff + retry.\n\n"
     "CPU-only inference: 2.46 s per explanation.\n"
     "  GPU deployment → < 0.5 s.\n\n"
     "Single-host only — no distributed\n"
     "or multi-network support yet.\n\n"
     "IPv6 not yet supported."),
    ("⚖","Ethical Considerations",
     "False positive risk:\n"
     "  Legitimate traffic could be blocked.\n"
     "  Mitigated by 60 s auto-revert.\n\n"
     "LLM explanations are AI-generated:\n"
     "  Should not replace human judgment\n"
     "  for critical security decisions.\n\n"
     "Balancing autonomous action with\n"
     "human oversight in production."),
    ("🌐","Scope Limitations",
     "Tested on local loopback only\n"
     "  (RFC 5737 TEST-NET source IPs).\n\n"
     "5 attack types only:\n"
     "  DNS amplification and HTTP\n"
     "  slowloris not yet covered.\n\n"
     "Adversarial disinformation evolves:\n"
     "  Thresholds need continuous\n"
     "  recalibration for new patterns."),
]
for i,(icon,title,body) in enumerate(three3):
    lx=Inches(0.42)+i*Inches(4.3)
    card(s,lx,Inches(1.28),Inches(4.1),Inches(6.0),
         title=title,icon=icon,lines=[(body,False)],ts=20,bs=13)

# ═══════════════════════════════════════════════════════════════════════════════
# S15  CONCLUSION & FUTURE SCOPE
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)
txt(s,"Conclusion & Future Scope",
    Inches(0.4),Inches(0.15),Inches(12.5),Inches(0.9),
    size=40,bold=True,color=CYAN_TITLE,align=PP_ALIGN.CENTER)
rect(s,Inches(0.4),Inches(1.06),Inches(12.5),Inches(0.04),fill=CYAN_LINE)
snum(s,15,TOTAL)

card(s,Inches(0.35),Inches(1.22),Inches(6.1),Inches(6.1),
     title="Summary",icon="🛡",
     lines=[
         ("Network Guardian proves that integrating local LLMs "
          "into network security is feasible, practical, and valuable.",False),
         ("",False),
         ("✓  Detects attacks in < 1 second",False),
         ("✓  Explains in plain English — 2.5 s",False),
         ("✓  Blocks in 0.17 s with 60 s auto-revert",False),
         ("✓  99% blocking success rate",False),
         ("✓  < 0.01 false positive rate — 100 hours tested",False),
         ("✓  99.96% Wireshark validation agreement",False),
         ("✓  100% user comprehension of AI explanations",False),
         ("   vs. < 30% for raw technical alerts",False),
         ("",False),
         ("Especially valuable as an educational tool:\n"
          "students launch attacks against themselves and\n"
          "immediately observe AI-driven detect+explain+block.",False),
     ],ts=19,bs=13.5)

card(s,Inches(6.7),Inches(1.22),Inches(6.3),Inches(6.1),
     title="Future Roadmap",icon="🚀",
     lines=[
         ("GPU-Accelerated LLM",True),
         ("  Llama-3 / Mistral on GPU → < 0.5 s analysis.",False),("",False),
         ("Adaptive Thresholds",True),
         ("  EWMA + Bayesian traffic baseline recalibration.",False),("",False),
         ("Multi-Host Deployment",True),
         ("  Distributed agents with centralised control.",False),("",False),
         ("SIEM Integration",True),
         ("  ThreatEvents → Elasticsearch / Splunk / syslog.",False),("",False),
         ("Threat Intelligence Feeds",True),
         ("  VirusTotal + MITRE ATT&CK + PostgreSQL.",False),("",False),
         ("Extended Attack Coverage",True),
         ("  DNS amplification, HTTP slowloris,\n"
          "  multiple LLM workers for parallelism.",False),
     ],ts=19,bs=13)

# ═══════════════════════════════════════════════════════════════════════════════
# S16  QUESTIONS  (Q&A closer)
# ═══════════════════════════════════════════════════════════════════════════════
s=prs.slides.add_slide(BLANK); bg(s)

# Neon border frame
rect(s,Inches(0.35),Inches(0.28),Inches(12.63),Inches(6.9),
     fill=BG,border=CYAN_GLOW,bw=Pt(2.5))

# Corner dots
for (lx,ty) in [(Inches(0.25),Inches(0.18)),(Inches(12.88),Inches(0.18)),
                 (Inches(0.25),Inches(7.08)),(Inches(12.88),Inches(7.08))]:
    rect(s,lx,ty,Inches(0.14),Inches(0.14),fill=CYAN_GLOW)

txt(s,"?",Inches(4.6),Inches(0.5),Inches(4.13),Inches(3.0),
    size=160,bold=True,color=CYAN_TITLE,align=PP_ALIGN.CENTER)

txt(s,"Questions?",Inches(0.5),Inches(3.65),Inches(12.3),Inches(1.0),
    size=54,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"Thank you for your attention.",
    Inches(0.5),Inches(4.72),Inches(12.3),Inches(0.58),
    size=24,color=LIGHT,align=PP_ALIGN.CENTER)

rect(s,Inches(2.5),Inches(5.4),Inches(8.33),Inches(0.04),fill=CYAN_LINE)

txt(s,"Arya Y P  |  Adi Narayan Prasad G  |  Ashrith S Jain  |  Abhishek",
    Inches(0.5),Inches(5.6),Inches(12.3),Inches(0.42),
    size=16,color=WHITE,align=PP_ALIGN.CENTER)
txt(s,"Dept. of ISE  |  B.M.S. College of Engineering, Bengaluru  |  ICWITE 2026",
    Inches(0.5),Inches(6.08),Inches(12.3),Inches(0.38),
    size=13,color=LIGHT,align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
