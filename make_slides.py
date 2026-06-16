"""
make_slides.py
Generates stunning 1920x1080 PNG slides for Network Guardian ICWITE 2026.
Each slide is a full rendered image — no plain shapes, every pixel designed.
"""

import os, math
import numpy as np
from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation
from pptx.util import Inches

W, H = 1920, 1080
SDIR  = r"c:\Users\adina\OneDrive\Desktop\cnfinal\slide_images"
CDIR  = r"c:\Users\adina\OneDrive\Desktop\cnfinal\ppt_images"
OUT   = r"c:\Users\adina\OneDrive\Desktop\cnfinal\NetworkGuardian_Conference.pptx"
FONTS = r"C:\Windows\Fonts"
os.makedirs(SDIR, exist_ok=True)

# ── Color Palette ─────────────────────────────────────────────────────────────
BG1  = (5,   10,  24)
BG2  = (8,   18,  44)
BG3  = (11,  26,  60)
CARD = (10,  22,  52)
CARD2= (7,   15,  36)
CARD3= (14,  30,  68)

CYN  = (0,   210, 255)   # Primary accent
CYN2 = (0,   140, 200)
CYN3 = (0,   60,  120)

PRP  = (140,  50, 255)   # Purple
PRP2 = (80,   20, 180)

GRN  = (0,   255, 120)   # Green
GRN2 = (0,   170,  80)

ORG  = (255, 140,   0)   # Orange
RED  = (255,  60,  60)

WHT  = (255, 255, 255)
LGR  = (180, 205, 225)   # Light gray
GRY  = (120, 148, 178)   # Gray
DGR  = (55,  75, 110)    # Dark gray

# ── Fonts ─────────────────────────────────────────────────────────────────────
_fc = {}
def F(sz, bold=False):
    k = (sz, bold)
    if k not in _fc:
        names = ['segoeuib.ttf','arialbd.ttf','calibrib.ttf'] if bold else \
                ['segoeui.ttf','arial.ttf','calibri.ttf']
        for n in names:
            p = os.path.join(FONTS, n)
            if os.path.exists(p):
                try:
                    _fc[k] = ImageFont.truetype(p, sz); break
                except: pass
        if k not in _fc:
            _fc[k] = ImageFont.load_default()
    return _fc[k]

# ── Background generators ─────────────────────────────────────────────────────
def make_bg(c1=BG1, c2=BG3):
    arr = np.zeros((H, W, 3), dtype=np.uint8)
    for y in range(H):
        t = y / H
        arr[y, :] = [int(c1[i]*(1-t)+c2[i]*t) for i in range(3)]
    return Image.fromarray(arr)

def add_dots(img, step=55, col=CYN3, alpha=22):
    ov = Image.new('RGBA', (W, H), (0,0,0,0))
    d  = ImageDraw.Draw(ov)
    for x in range(0, W, step):
        for y in range(0, H, step):
            d.ellipse([x-1,y-1,x+2,y+2], fill=(*col, alpha))
    return Image.alpha_composite(img.convert('RGBA'), ov).convert('RGB')

def add_grid(img, step=90, col=CYN3, alpha=14):
    ov = Image.new('RGBA', (W, H), (0,0,0,0))
    d  = ImageDraw.Draw(ov)
    for x in range(0, W, step):
        d.line([x,0,x,H], fill=(*col, alpha), width=1)
    for y in range(0, H, step):
        d.line([0,y,W,y], fill=(*col, alpha), width=1)
    return Image.alpha_composite(img.convert('RGBA'), ov).convert('RGB')

def add_orb(img, cx, cy, r, col, intensity=70):
    ov = Image.new('RGBA', (W, H), (0,0,0,0))
    d  = ImageDraw.Draw(ov)
    step = max(1, r//20)
    for ri in range(r, 0, -step):
        a = int(intensity * (1 - ri/r)**1.8)
        d.ellipse([cx-ri, cy-ri, cx+ri, cy+ri], fill=(*col, min(a,255)))
    ov = ov.filter(ImageFilter.GaussianBlur(r//5))
    return Image.alpha_composite(img.convert('RGBA'), ov).convert('RGB')

def add_network(img, nodes, col=CYN3, nc=CYN, alpha_line=20, alpha_node=160):
    ov = Image.new('RGBA', (W, H), (0,0,0,0))
    d  = ImageDraw.Draw(ov)
    for i, (x1,y1) in enumerate(nodes):
        for j, (x2,y2) in enumerate(nodes):
            if i < j:
                dist = math.hypot(x2-x1, y2-y1)
                if dist < 500:
                    a = int(alpha_line * (1 - dist/500))
                    d.line([x1,y1,x2,y2], fill=(*col, a), width=1)
    for nx, ny in nodes:
        d.ellipse([nx-4,ny-4,nx+4,ny+4], fill=(*nc, alpha_node))
        d.ellipse([nx-10,ny-10,nx+10,ny+10], outline=(*nc, 60), width=1)
    return Image.alpha_composite(img.convert('RGBA'), ov).convert('RGB')

# ── Drawing primitives ────────────────────────────────────────────────────────
def card(img, x, y, w, h, fill=CARD, border=CYN, brad=18, bw=2, glow_=True):
    ov = Image.new('RGBA', (W, H), (0,0,0,0))
    d  = ImageDraw.Draw(ov)
    d.rounded_rectangle([x,y,x+w,y+h], radius=brad, fill=(*fill, 245))
    if glow_:
        for i in range(12, 0, -2):
            a = int(45 * i/12)
            d.rounded_rectangle([x-i,y-i,x+w+i,y+h+i], radius=brad+i,
                                 outline=(*border, a), width=1)
    d.rounded_rectangle([x,y,x+w,y+h], radius=brad, outline=(*border, 220), width=bw)
    return Image.alpha_composite(img.convert('RGBA'), ov).convert('RGB')

def glow(img, text, xy, font, col, gcol=None, gr=22, anchor='lt'):
    if gcol is None: gcol = col
    ov = Image.new('RGBA', (W, H), (0,0,0,0))
    d  = ImageDraw.Draw(ov)
    d.text(xy, text, font=font, fill=(*gcol, 220), anchor=anchor)
    blurred = ov.filter(ImageFilter.GaussianBlur(gr))
    result  = Image.alpha_composite(img.convert('RGBA'), blurred)
    ImageDraw.Draw(result).text(xy, text, font=font, fill=(*col, 255), anchor=anchor)
    return result.convert('RGB')

def hline(img, y, x1=80, x2=None, col=CYN, w=2, alpha=180):
    if x2 is None: x2 = W-80
    ov = Image.new('RGBA', (W, H), (0,0,0,0))
    d  = ImageDraw.Draw(ov)
    d.line([x1,y,x2,y], fill=(*col, 50), width=8)
    d.line([x1,y,x2,y], fill=(*col, alpha), width=w)
    return Image.alpha_composite(img.convert('RGBA'), ov).convert('RGB')

def txtw(img_or_draw, text, x, y, font, col, max_w=None, line_gap=10, anchor='lt'):
    if isinstance(img_or_draw, ImageDraw.ImageDraw):
        draw = img_or_draw; ret_img = None
    else:
        ret_img = img_or_draw; draw = ImageDraw.Draw(img_or_draw)

    def render(draw_, text_, x_, y_):
        if anchor == 'mt':
            bb = draw_.textbbox((0,0), text_, font=font)
            x_ = x_ - (bb[2]-bb[0])//2
        draw_.text((x_,y_), text_, font=font, fill=col, anchor='lt')
        bb = draw_.textbbox((0,0), text_, font=font)
        return bb[3]-bb[1]

    if max_w is None:
        h_ = render(draw, text, x, y)
        return (ret_img, y+h_) if ret_img else y+h_

    words = text.split()
    lines_, cur = [], []
    for w_ in words:
        test = ' '.join(cur+[w_])
        bb = draw.textbbox((0,0), test, font=font)
        if (bb[2]-bb[0]) > max_w and cur:
            lines_.append(' '.join(cur)); cur = [w_]
        else:
            cur.append(w_)
    if cur: lines_.append(' '.join(cur))

    cy = y
    for line in lines_:
        h_ = render(draw, line, x, cy)
        cy += h_ + line_gap
    return (ret_img, cy) if ret_img else cy

def ctxt(draw, text, cy, font, col):
    bb = draw.textbbox((0,0), text, font=font)
    x  = (W-(bb[2]-bb[0]))//2
    draw.text((x, cy), text, font=font, fill=col)
    return cy + (bb[3]-bb[1])

def tw(draw, text, font):
    bb = draw.textbbox((0,0), text, font=font)
    return bb[2]-bb[0]

def th(draw, text, font):
    bb = draw.textbbox((0,0), text, font=font)
    return bb[3]-bb[1]

def header(img, title, subtitle=None):
    img = glow(img, title, (W//2, 42), F(72,True), WHT, CYN, 28, 'mt')
    img = hline(img, 135)
    if subtitle:
        d = ImageDraw.Draw(img)
        ctxt(d, subtitle, 148, F(28), GRY)
    return img

# ── Icon helpers (PIL-drawn, no emoji needed) ─────────────────────────────────
def draw_circle_badge(img, cx, cy, r, col, letter, lsz=None):
    """Draw a glowing circle badge with a letter inside."""
    ov = Image.new('RGBA',(W,H),(0,0,0,0))
    d  = ImageDraw.Draw(ov)
    d.ellipse([cx-r,cy-r,cx+r,cy+r], fill=(*col,25))
    for i in range(8,0,-2):
        d.ellipse([cx-r-i,cy-r-i,cx+r+i,cy+r+i], outline=(*col,int(35*i/8)), width=1)
    d.ellipse([cx-r,cy-r,cx+r,cy+r], outline=(*col,200), width=3)
    img = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')
    d2  = ImageDraw.Draw(img)
    sz  = lsz or int(r*1.1)
    bb  = d2.textbbox((0,0),letter,font=F(sz,True))
    lw,lh = bb[2]-bb[0], bb[3]-bb[1]
    d2.text((cx-lw//2, cy-lh//2-2), letter, font=F(sz,True), fill=col)
    return img

def draw_check(draw, cx, cy, col, size=22):
    """Draw a checkmark."""
    pts = [(cx-size//2, cy+2), (cx-size//8, cy+size//2), (cx+size//2, cy-size//3)]
    for i in range(len(pts)-1):
        draw.line([pts[i], pts[i+1]], fill=col, width=4)

def draw_xmark(draw, cx, cy, col, size=18):
    """Draw an X mark."""
    s = size//2
    draw.line([cx-s,cy-s,cx+s,cy+s], fill=col, width=4)
    draw.line([cx+s,cy-s,cx-s,cy+s], fill=col, width=4)

def draw_dot(img, cx, cy, r, col):
    """Draw a glowing filled dot."""
    ov = Image.new('RGBA',(W,H),(0,0,0,0))
    d  = ImageDraw.Draw(ov)
    d.ellipse([cx-r,cy-r,cx+r,cy+r], fill=(*col,220))
    d.ellipse([cx-r-3,cy-r-3,cx+r+3,cy+r+3], outline=(*col,80), width=2)
    return Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')

def snum(img, n, t=14):
    d = ImageDraw.Draw(img)
    d.text((W-80, H-42), f"{n} / {t}", font=F(22), fill=DGR, anchor='rt')
    return img

def paste(img, fname, x, y, w, h):
    p = os.path.join(CDIR, fname)
    if not os.path.exists(p): return img
    ch = Image.open(p).convert('RGBA').resize((w,h), Image.LANCZOS)
    tmp = img.convert('RGBA')
    tmp.paste(ch, (x,y), ch)
    return tmp.convert('RGB')

def save_s(img, n):
    p = os.path.join(SDIR, f"slide_{n:02d}.png")
    img.save(p, 'PNG')
    print(f"  [{n:02d}] saved")
    return p

PATHS = []
TOTAL = 14

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
def s01():
    img = make_bg(BG1, (10,22,55))
    img = add_dots(img, 52, CYN3, 20)

    # Background orbs
    img = add_orb(img, 250, 200, 480, (0,50,120), 55)
    img = add_orb(img, W-200, H-150, 380, (70,20,160), 45)

    # Network nodes decoration
    np.random.seed(42)
    nodes = [(np.random.randint(80,W-80), np.random.randint(80,H-80)) for _ in range(22)]
    img = add_network(img, nodes, CYN3, CYN, 18, 150)

    # Top accent bar
    ov = Image.new('RGBA',(W,H),(0,0,0,0))
    d  = ImageDraw.Draw(ov)
    for i,col in enumerate([CYN,PRP,GRN]):
        d.rectangle([i*4, 0, i*4+3, 6], fill=(*col,255))
    d.rectangle([0,0,W,5], fill=(*CYN,80))
    img = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')

    # ── Main title
    img = glow(img, "NETWORK GUARDIAN", (W//2, 265), F(116,True), WHT, CYN, 40, 'mm')

    # Underline
    img = hline(img, 330, W//2-460, W//2+460, CYN, 2, 160)

    d = ImageDraw.Draw(img)

    # Subtitle
    ctxt(d, "An Intelligent System That Automatically Protects Computer Networks", 348, F(32), LGR)
    ctxt(d, "from Attacks Using Artificial Intelligence", 390, F(32), LGR)

    # Authors strip
    img = card(img, 380, 445, W-760, 64, fill=(10,20,50), border=CYN, brad=32, bw=1)
    d = ImageDraw.Draw(img)
    ctxt(d, "Arya Y P   ·   Adi Narayan Prasad G   ·   Ashrith S Jain   ·   Abhishek", 460, F(26,True), CYN)
    ctxt(d, "Guide: Prof. Nalina V   |   Dept. of ISE   |   BMSCE, Bengaluru", 506, F(24), GRY)

    # Conference badge
    img = card(img, W//2-185, 585, 370, 56, fill=CARD2, border=PRP, brad=28, bw=2)
    d = ImageDraw.Draw(img)
    ctxt(d, "ICWITE 2026  —  ACCEPTED PAPER", 603, F(24,True), PRP)

    # 3 metric boxes
    metrics = [("<1 SEC", "Attack Detection"), ("2.88 SEC", "End-to-End Response"), ("99%", "Blocking Success")]
    bw2 = 340
    gap  = 36
    sx   = (W - (bw2*3 + gap*2))//2
    for i,(val,lbl) in enumerate(metrics):
        bx = sx + i*(bw2+gap)
        img = card(img, bx, 680, bw2, 120, fill=CARD2, border=CYN, brad=16)
        d   = ImageDraw.Draw(img)
        vw  = tw(d, val, F(52,True))
        d.text((bx+(bw2-vw)//2, 696), val, font=F(52,True), fill=CYN)
        lw  = tw(d, lbl, F(22))
        d.text((bx+(bw2-lw)//2, 756), lbl, font=F(22), fill=LGR)

    # Bottom bar
    img = hline(img, H-6, 0, W, CYN, 4, 100)
    return save_s(img, 1)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═════════════════════════════════════════════════════════════════════════════
def s02():
    img = make_bg(BG1, (20,8,8))
    img = add_dots(img, 60, (80,20,20), 18)
    img = add_orb(img, W//2, H//2, 700, (60,10,10), 40)

    img = header(img, "The Cybersecurity Crisis")
    img = snum(img, 2)

    d = ImageDraw.Draw(img)

    # Big stat
    img = glow(img, "207", (W//2, 320), F(280,True), (255,60,60), (180,0,0), 50, 'mm')
    d   = ImageDraw.Draw(img)
    ctxt(d, "DAYS  —  average time to detect a data breach  [IBM 2023]", 480, F(30), LGR)

    img = hline(img, 540, 200, W-200, (180,40,40), 1, 120)
    d   = ImageDraw.Draw(img)

    # 4 problem cards
    probs = [
        ("CRYPTIC ALERTS",
         "IDS shows: \"ICMP flood from 192.168.1.50\"\nOperators cannot act without expert training."),
        ("ALERT FATIGUE",
         "Thousands of alerts daily.\nSecurity teams miss critical threats."),
        ("MANUAL RESPONSE",
         "Attacks complete in milliseconds.\n+73 days to contain after detection."),
        ("SKILLS GAP",
         "Shortage of trained security personnel\nin most organizations worldwide."),
    ]
    cw, ch_ = 420, 230
    gap2     = 25
    sx2      = (W - (cw*4 + gap2*3))//2
    for i,(title,body) in enumerate(probs):
        bx = sx2 + i*(cw+gap2)
        img = card(img, bx, 568, cw, ch_, fill=(18,8,8), border=(200,50,50), brad=16)
        img = draw_dot(img, bx+22, 594, 9, RED)
        d   = ImageDraw.Draw(img)
        d.text((bx+42, 582), title, font=F(24,True), fill=(255,90,90))
        txtw(d, body, bx+20, 636, F(22), LGR, cw-40)

    img = hline(img, 830, 200, W-200, (180,40,40), 1, 80)
    d   = ImageDraw.Draw(img)
    ctxt(d, "We need an intelligent system that detects, explains, and responds — autonomously.", 848, F(28,True), WHT)

    return save_s(img, 2)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION (3 pillars)
# ═════════════════════════════════════════════════════════════════════════════
def s03():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, 0,   H//2, 400, CYN3, 50)
    img = add_orb(img, W//2, H-50, 350, PRP2,  40)
    img = add_orb(img, W, H//2,  350, GRN2,  35)

    img = header(img, "Our Solution: Network Guardian",
                 "Three core capabilities working in concert")
    img = snum(img, 3)

    pillars = [
        (CYN,  "D",  "DETECT",
         "Sub-second attack identification",
         ["Sliding-window counters per source IP",
          "5 attack types monitored in parallel",
          "< 1 second detection latency",
          "< 0.01 false positive rate (100 hrs)"]),
        (PRP,  "AI", "EXPLAIN",
         "Plain English AI explanations",
         ["TinyLlama 1.1B via Ollama (local)",
          "No cloud — private & offline",
          "94% first-attempt JSON success",
          "Mean analysis: 2.46 seconds"]),
        (GRN,  "S",  "PROTECT",
         "Autonomous OS-level mitigation",
         ["route blackhole blocks attacker IP",
          "60-second auto-revert daemon",
          "99% blocking success (198/200)",
          "100% traffic reduction verified"]),
    ]

    cw3, ch3 = 540, 490
    gap3 = 40
    sx3  = (W - (cw3*3 + gap3*2))//2

    for i,(col, badge_letter, title, subtitle, points) in enumerate(pillars):
        bx = sx3 + i*(cw3+gap3)
        img = card(img, bx, 185, cw3, ch3, fill=CARD2, border=col, brad=20, bw=3)

        # Top color strip
        ov = Image.new('RGBA',(W,H),(0,0,0,0))
        d_ = ImageDraw.Draw(ov)
        d_.rounded_rectangle([bx,185,bx+cw3,228], radius=20, fill=(*col,40))
        img = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')

        # Circle badge with letter
        img = draw_circle_badge(img, bx+cw3//2, 236, 42, col, badge_letter, 38)

        # Title with glow
        img = glow(img, title, (bx+cw3//2, 288), F(44,True), col, col, 18, 'mt')
        d   = ImageDraw.Draw(img)

        # Subtitle
        stw = tw(d, subtitle, F(23))
        d.text((bx+(cw3-stw)//2, 342), subtitle, font=F(23), fill=LGR)

        # Divider
        d.line([bx+30, 376, bx+cw3-30, 376], fill=(*col, 80), width=1)

        # Bullet points
        by = 392
        for pt in points:
            d.text((bx+30, by), ">", font=F(25,True), fill=col)
            txtw(d, pt, bx+58, by+2, F(24), WHT, cw3-88)
            by += 52

    return save_s(img, 3)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SYSTEM ARCHITECTURE
# ═════════════════════════════════════════════════════════════════════════════
def s04():
    img = make_bg()
    img = add_grid(img, 80, CYN3, 12)
    img = add_orb(img, W-180, 200, 350, PRP2, 35)

    img = header(img, "System Architecture",
                 "5-Layer modular queue-based pipeline — each layer independently replaceable")
    img = snum(img, 4)

    layers = [
        (CYN,  "LAYER 5", "Web Dashboard",   "Flask + Socket.IO  |  Browser UI  |  REST API  |  500ms push"),
        (PRP,  "LAYER 4", "Mitigator",        "OS route blackhole  |  60s auto-revert  |  shell=False  |  IP validation"),
        (GRN,  "LAYER 3", "LLM Analyzer",     "TinyLlama 1.1B via Ollama  |  JSON output  |  Retry policy  |  Fallback"),
        (ORG,  "LAYER 2", "Attack Detector",  "5 sliding-window detectors  |  SlidingWindowCounter  |  PortScanTracker"),
        (CYN,  "LAYER 1", "Network Sniffer",  "Scapy packet capture  |  Npcap (Win) / libpcap (Linux)  |  Metadata extraction"),
    ]

    lh  = 120
    gap4 = 8
    sx4  = 120
    lw4  = 1010
    sy4  = 175

    for i,(col, lnum, lname, ldesc) in enumerate(layers):
        by4 = sy4 + i*(lh+gap4)
        img = card(img, sx4, by4, lw4, lh, fill=CARD2, border=col, brad=14, bw=2)

        # Left accent strip
        ov = Image.new('RGBA',(W,H),(0,0,0,0))
        d_ = ImageDraw.Draw(ov)
        d_.rounded_rectangle([sx4, by4, sx4+8, by4+lh], radius=7, fill=(*col,220))
        img = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')

        d = ImageDraw.Draw(img)
        d.text((sx4+24, by4+10), lnum,   font=F(20,True), fill=col)
        d.text((sx4+24, by4+34), lname,  font=F(32,True), fill=WHT)
        d.text((sx4+24, by4+76), ldesc,  font=F(20),      fill=GRY)

        # Arrow down (not on last)
        if i < len(layers)-1:
            ax = sx4 + lw4//2
            ay = by4 + lh + 2
            d.polygon([(ax-12,ay+1),(ax+12,ay+1),(ax,ay+gap4-1)], fill=(*col,150))

    # Right side: queue bus diagram
    rx = 1200
    qitems = [
        ("packet_queue",  "cap 5,000",  CYN),
        ("threat_queue",  "cap 100",    ORG),
        ("ui_log_queue",  "cap 500",    PRP),
        ("active_blocks", "dict + lock",GRN),
        ("threat_history","last 10",    GRY),
    ]
    img = card(img, rx, 185, 640, 680, fill=CARD2, border=DGR, brad=16, bw=1, glow_=False)
    d   = ImageDraw.Draw(img)
    d.text((rx+20, 198), "Shared State Bus", font=F(30,True), fill=CYN)
    d.line([rx+20, 236, rx+620, 236], fill=(*CYN,60), width=1)
    for j,(name,cap,col) in enumerate(qitems):
        qy = 250 + j*100
        img = card(img, rx+20, qy, 600, 82, fill=CARD, border=col, brad=12, bw=2)
        d   = ImageDraw.Draw(img)
        d.text((rx+40, qy+12), name, font=F(28,True), fill=col)
        d.text((rx+40, qy+48), cap,  font=F(22), fill=GRY)

    # Key stat
    img = card(img, rx, 890, 640, 100, fill=(5,30,10), border=GRN, brad=14, bw=2)
    d   = ImageDraw.Draw(img)
    d.text((rx+20, 903), "5-second buffer @ 1,000 pps", font=F(28,True), fill=GRN)
    d.text((rx+20, 940), "Single-responsibility · Fault isolated · Queue-based", font=F(22), fill=LGR)

    return save_s(img, 4)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DETECTION ENGINE
# ═════════════════════════════════════════════════════════════════════════════
def s05():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, 300, H//2, 400, CYN3, 45)

    img = header(img, "Attack Detection Engine",
                 "Five sliding-window detectors — independent, parallel, sub-second")
    img = snum(img, 5)

    # Embed radar chart
    img = paste(img, 'img_radar.png', 80, 168, 820, 820)

    # Right: table + latency
    attacks = [
        ("ICMP Flood",       "20 pps",  "1.0 s", "0.14 s", CYN),
        ("SYN Flood",        "20 pps",  "1.0 s", "0.27 s", PRP),
        ("UDP Flood",        "20 pps",  "1.0 s", "0.15 s", ORG),
        ("Port Scan",        "10 ports","5.0 s", "0.51 s", GRN),
        ("IP Fragmentation", "10 pps",  "1.0 s", "0.21 s", (255,180,0)),
    ]

    tx_ = 980
    d   = ImageDraw.Draw(img)
    d.text((tx_, 178), "Detection Parameters  &  Latency", font=F(34,True), fill=CYN)

    # Table headers
    col_x = [tx_, tx_+270, tx_+440, tx_+590]
    col_h2 = ["Attack Type", "Threshold", "Window", "Mean Latency"]
    hrow_y = 222
    for ci,(hd,cx) in enumerate(zip(col_h2,col_x)):
        d.text((cx, hrow_y), hd, font=F(23,True), fill=GRY)
    d.line([tx_, 250, W-60, 250], fill=(*CYN3,150), width=1)

    for ri,(name,thresh,win,lat,col) in enumerate(attacks):
        ry_ = 260 + ri*78
        img = card(img, tx_-10, ry_-8, W-tx_-50, 68,
                   fill=CARD if ri%2==0 else CARD2, border=col, brad=10, bw=1, glow_=False)
        d   = ImageDraw.Draw(img)
        # Color dot
        d.ellipse([tx_+2, ry_+18, tx_+18, ry_+34], fill=col)
        d.text((tx_+26,        ry_+8),  name,   font=F(25,True), fill=WHT)
        d.text((col_x[1],      ry_+8),  thresh,  font=F(25), fill=LGR)
        d.text((col_x[2],      ry_+8),  win,     font=F(25), fill=LGR)
        # Latency with bar
        d.text((col_x[3],      ry_+8),  lat,     font=F(25,True), fill=col)
        # Mini latency bar
        blen = int(float(lat.replace(' s','')) / 0.62 * 140)
        ov2 = Image.new('RGBA',(W,H),(0,0,0,0))
        d2  = ImageDraw.Draw(ov2)
        d2.rounded_rectangle([col_x[3]+85, ry_+16, col_x[3]+85+blen, ry_+28],
                              radius=4, fill=(*col, 180))
        img = Image.alpha_composite(img.convert('RGBA'),ov2).convert('RGB')
        d   = ImageDraw.Draw(img)

    # Key code snippet
    img = card(img, tx_-10, 665, W-tx_-50, 220, fill=(4,14,32), border=DGR, brad=14, bw=1, glow_=False)
    d   = ImageDraw.Draw(img)
    d.text((tx_+10, 675), "Core Algorithm  (SlidingWindowCounter)", font=F(24,True), fill=GRY)
    code = [
        ("class ", WHT), ("SlidingWindowCounter", CYN), (":", WHT),
    ]
    d.text((tx_+10, 710), "class SlidingWindowCounter:", font=F(22), fill=(140,200,255))
    d.text((tx_+10, 738), "  deque of timestamps + lock per source IP", font=F(22), fill=LGR)
    d.text((tx_+10, 765), "  prune entries older than window size", font=F(22), fill=LGR)
    d.text((tx_+10, 792), "  if count > threshold: raise ThreatEvent", font=F(22,True), fill=GRN)
    d.text((tx_+10, 820), "  60s cooldown per (IP, attack_type) pair", font=F(22), fill=GRY)
    d.text((tx_+10, 850), "< 0.01 false positive rate  ·  100 hours tested", font=F(24,True), fill=CYN)

    return save_s(img, 5)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI EXPLANATION (Before vs After)
# ═════════════════════════════════════════════════════════════════════════════
def s06():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, W//2, H//2, 600, PRP2, 35)

    img = header(img, "AI-Powered Explanation Layer",
                 "TinyLlama 1.1B via Ollama  ·  Local & Private  ·  No Cloud Required")
    img = snum(img, 6)

    # LEFT: Terminal - raw alert  (width=740, ends at x=800)
    img = card(img, 55, 188, 740, 760, fill=(4,10,4), border=(80,200,80), brad=16, bw=2)
    ov  = Image.new('RGBA',(W,H),(0,0,0,0))
    d_  = ImageDraw.Draw(ov)
    d_.rounded_rectangle([55,188,795,210], radius=16, fill=(20,80,20,200))
    img = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')
    img = draw_dot(img, 75, 200, 7, GRN)
    d   = ImageDraw.Draw(img)
    d.text((92, 192), "Traditional IDS Alert", font=F(24,True), fill=(100,220,100))
    d.text((80, 238), "$ snort -A console -c snort.conf", font=F(22), fill=(80,180,80))
    d.text((80, 268), "", font=F(22), fill=WHT)
    alert_lines = [
        ("[**] [1:1000001:1] ICMP Flood Detected [**]",     (255,80,80)),
        ("Classification: Denial of Service Attack",          (200,80,80)),
        ("Priority: 1",                                       (200,80,80)),
        ("",                                                   WHT),
        ("06/15-14:23:17.482156 203.0.113.50:0 ->",          (180,220,180)),
        ("192.168.1.100:0",                                   (180,220,180)),
        ("PROTO:1 TOS:0x0 TTL:64 ID:1234 IpLen:20 DgmLen:84",(140,180,140)),
        ("Type:8 Code:0 ID:10 Seq:4150 ECHO",                (140,180,140)),
        ("",                                                   WHT),
        ("0x0000: 45 00 00 54 04 D2 00 00  40 01 F6 4E C0 A8", (100,160,100)),
        ("0x0010: 01 64 C0 A8 01 64 08 00  EC 4B 00 0A 10 36", (100,160,100)),
        ("",                                                   WHT),
        ("!  Non-expert cannot interpret this output",         (255,200,80)),
        ("!  No automatic action taken",                       (255,200,80)),
        ("!  Manual intervention required",                    (255,200,80)),
    ]
    ay = 236
    for line_,col in alert_lines:
        d.text((72, ay), line_, font=F(19), fill=col)
        ay += 28

    # CENTER arrow — drawn with PIL lines, no emoji
    cx6 = 957  # centre of 795..1120 gap
    # Draw arrow shaft and head
    ov_arr = Image.new('RGBA',(W,H),(0,0,0,0))
    da = ImageDraw.Draw(ov_arr)
    da.line([860, H//2, 1050, H//2], fill=(*CYN,200), width=4)
    da.polygon([(1050,H//2-14),(1050,H//2+14),(1085,H//2)], fill=(*CYN,220))
    blurred_arr = ov_arr.filter(ImageFilter.GaussianBlur(6))
    img = Image.alpha_composite(img.convert('RGBA'), blurred_arr).convert('RGB')
    da2 = ImageDraw.Draw(img)
    da2.line([860, H//2, 1050, H//2], fill=CYN, width=3)
    da2.polygon([(1050,H//2-12),(1050,H//2+12),(1082,H//2)], fill=CYN)
    nw = tw(da2, "Network", F(22,True))
    da2.text((cx6-nw//2, H//2-52), "Network", font=F(22,True), fill=CYN)
    gw = tw(da2, "Guardian", F(22,True))
    da2.text((cx6-gw//2, H//2-24), "Guardian", font=F(22,True), fill=CYN)

    # RIGHT: AI explanation (starts at x=1120, width=750, ends at 1870)
    img = card(img, 1120, 188, 750, 760, fill=CARD2, border=CYN, brad=16, bw=2)
    ov2  = Image.new('RGBA',(W,H),(0,0,0,0))
    d2_  = ImageDraw.Draw(ov2)
    d2_.rounded_rectangle([1120,188,1870,218], radius=16, fill=(*CYN3,200))
    img  = Image.alpha_composite(img.convert('RGBA'),ov2).convert('RGB')
    img  = draw_dot(img, 1142, 202, 7, CYN)
    d    = ImageDraw.Draw(img)
    d.text((1158, 192), "Network Guardian AI Analysis", font=F(24,True), fill=CYN)

    # Right panel uses x=1140 (inside the 1120-wide card), max width ~710
    RP = 1140   # right panel text left edge
    sections = [
        ("ATTACK DETECTED",    "SYN Flood",                          RED,  CYN),
        ("SOURCE IP",          "203.0.113.50",                        GRY,  WHT),
        ("SEVERITY",           "HIGH  87 pps  (threshold: 20 pps)",   GRY,  ORG),
        ("", "", GRY, WHT),
        ("PLAIN-ENGLISH EXPLANATION", "", GRY, GRY),
    ]
    sy6 = 236
    for lbl,val,lc,vc in sections:
        if lbl == "PLAIN-ENGLISH EXPLANATION":
            d.text((RP, sy6), lbl+":", font=F(20,True), fill=GRY)
            sy6 += 32
            explanation = (
                "The attacker sends TCP SYN packets without completing\n"
                "the three-way handshake. This exhausts the server's\n"
                "TCP connection queue (128-512 slots), preventing\n"
                "legitimate connections. Classic Denial-of-Service."
            )
            for line_ in explanation.split('\n'):
                d.text((RP, sy6), line_, font=F(21), fill=LGR)
                sy6 += 28
            sy6 += 8
        elif lbl:
            d.text((RP,    sy6), lbl+":", font=F(20,True), fill=lc)
            d.text((RP+220,sy6), val,     font=F(21,True), fill=vc)
            sy6 += 34

    d.line([RP, sy6+4, 1855, sy6+4], fill=(*CYN,60), width=1)
    sy6 += 16

    d.text((RP, sy6), "RECOMMENDED ACTION:", font=F(20,True), fill=GRY)
    sy6 += 30
    d.text((RP, sy6), "Block IP 203.0.113.50 immediately", font=F(23,True), fill=GRN)
    sy6 += 34
    d.text((RP, sy6), "route add 203.0.113.50 mask 255.255.255.255 192.0.2.1", font=F(19), fill=(100,200,255))
    sy6 += 32
    d.text((RP, sy6), "Auto-revert in 60 seconds (prevents permanent block)", font=F(19), fill=GRY)
    sy6 += 42

    # Execution badge
    img = card(img, RP-10, sy6, 720, 76, fill=(5,25,10), border=GRN, brad=12, bw=2)
    d   = ImageDraw.Draw(img)
    draw_check(d, RP+20, sy6+22, GRN, 20)
    d.text((RP+46, sy6+6),  "MITIGATION APPLIED",                       font=F(26,True), fill=GRN)
    d.text((RP+10, sy6+40), "203.0.113.50 --> NULL routed  |  0.17 s",  font=F(20),      fill=LGR)

    sy6 += 92
    # Stats row — 4 equal boxes
    stats6 = [("94%","JSON Success"), ("2.46s","Mean Analysis"), ("3 Retries","Policy"), ("0.1","Temp")]
    sw6 = 174
    for j,(v,l) in enumerate(stats6):
        bx6 = RP-10 + j*(sw6+6)
        img = card(img, bx6, sy6, sw6, 68, fill=CARD, border=PRP, brad=10, bw=1)
        d   = ImageDraw.Draw(img)
        vw6 = tw(d, v, F(24,True))
        d.text((bx6+(sw6-vw6)//2, sy6+6),  v, font=F(24,True), fill=PRP)
        lw6 = tw(d, l, F(18))
        d.text((bx6+(sw6-lw6)//2, sy6+38), l, font=F(18), fill=GRY)

    return save_s(img, 6)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AUTONOMOUS MITIGATION
# ═════════════════════════════════════════════════════════════════════════════
def s07():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, W//2, 300, 500, (0,60,20), 40)

    img = header(img, "Autonomous Mitigation",
                 "OS-level route blocking  ·  Injection-proof  ·  Auto-revert in 60 seconds")
    img = snum(img, 7)

    # Timeline visualization
    timeline = [
        (CYN,  "0.00s", "ATTACK\nDETECTED",   "Sliding-window\nthreshold exceeded"),
        (PRP,  "2.46s", "AI ANALYSIS\nCOMPLETE","Plain English\nexplanation ready"),
        (ORG,  "2.63s", "ROUTE\nBLOCKED",      "OS-level null\nroute applied"),
        (GRN,  "62.6s", "AUTO\nREVERT",         "Route deleted\nauto by daemon"),
    ]

    # Timeline bar
    tx1, tx2 = 120, W-120
    ty_line  = 440
    ov   = Image.new('RGBA',(W,H),(0,0,0,0))
    d_   = ImageDraw.Draw(ov)
    d_.line([tx1, ty_line, tx2, ty_line], fill=(*CYN3,150), width=4)
    img  = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')

    # Spread positions evenly to prevent card overlap  (card width=340)
    positions = [0.0, 0.30, 0.62, 1.0]
    for i,((col,time,title,desc),pos) in enumerate(zip(timeline,positions)):
        px = int(tx1 + pos*(tx2-tx1))

        # Timeline dot
        ov2 = Image.new('RGBA',(W,H),(0,0,0,0))
        d2  = ImageDraw.Draw(ov2)
        d2.ellipse([px-14, ty_line-14, px+14, ty_line+14], fill=(*col,220))
        d2.ellipse([px-22, ty_line-22, px+22, ty_line+22], outline=(*col,80), width=2)
        img = Image.alpha_composite(img.convert('RGBA'),ov2).convert('RGB')

        # Time label above
        d = ImageDraw.Draw(img)
        tw_ = tw(d, time, F(26,True))
        d.text((px-tw_//2, ty_line-70), time, font=F(26,True), fill=col)

        # Card below — width=340, clamp to slide edges
        cw_t = 340; ch_t = 190
        cx_t = max(60, min(px-cw_t//2, W-cw_t-60))
        img  = card(img, cx_t, ty_line+30, cw_t, ch_t, fill=CARD2, border=col, brad=14, bw=2)
        d    = ImageDraw.Draw(img)
        # Title in card
        for j,line_ in enumerate(title.split('\n')):
            lw_ = tw(d, line_, F(26,True))
            d.text((cx_t+(cw_t-lw_)//2, ty_line+46+j*34), line_, font=F(26,True), fill=col)
        # Desc
        for j,line_ in enumerate(desc.split('\n')):
            lw_ = tw(d, line_, F(22))
            d.text((cx_t+(cw_t-lw_)//2, ty_line+122+j*30), line_, font=F(22), fill=LGR)

    # Bottom: 3 cards for safety
    safety = [
        (CYN, "Command Injection",
         "LLM output NEVER passed to shell.\n"
         "IP validated via Python ipaddress\n"
         "module before any OS call.\n"
         "subprocess.run(shell=False)"),
        (GRN, "Auto-Revert Daemon",
         "Daemon thread sleeps exactly 60s\n"
         "then runs route delete command.\n"
         "False positives self-heal.\n"
         "No permanent routing damage."),
        (ORG, "Duplicate Prevention",
         "IP marked in active_blocks dict\n"
         "before route command executes.\n"
         "Loopback and multicast refused.\n"
         "Lock protects concurrent access."),
    ]
    cw_s, ch_s = 550, 200
    gap_s = 40
    sx_s  = (W - (cw_s*3 + gap_s*2))//2
    for i,(col,title,body) in enumerate(safety):
        bx_s = sx_s + i*(cw_s+gap_s)
        img  = card(img, bx_s, 720, cw_s, ch_s, fill=CARD2, border=col, brad=16, bw=2)
        d    = ImageDraw.Draw(img)
        d.text((bx_s+20, 736), title, font=F(26,True), fill=col)
        txtw(d, body, bx_s+20, 780, F(22), LGR, cw_s-40)

    return save_s(img, 7)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS (Big Numbers)
# ═════════════════════════════════════════════════════════════════════════════
def s08():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, W//2, 250, 600, (0,60,30), 45)

    img = header(img, "Experimental Results",
                 "Windows 11  ·  Npcap loopback  ·  TinyLlama CPU inference  ·  200 test attacks")
    img = snum(img, 8)

    # 4 BIG number cards
    big = [
        (CYN,  "<1",   "SECOND",     "Attack detection\nlatency"),
        (PRP,  "2.88", "SECONDS",    "Total pipeline\nend-to-end"),
        (GRN,  "99%",  "SUCCESS",    "Blocking success\n198 of 200 tests"),
        (ORG,  "<0.01","FALSE +VE",  "False positive\nrate, 100 hours"),
    ]
    bw8 = 400; bh8 = 250
    gap8 = 24
    sx8  = (W - (bw8*4 + gap8*3))//2

    for i,(col,num,unit,desc) in enumerate(big):
        bx8 = sx8 + i*(bw8+gap8)
        img = card(img, bx8, 185, bw8, bh8, fill=CARD2, border=col, brad=18, bw=3)

        # Top glow strip
        ov = Image.new('RGBA',(W,H),(0,0,0,0))
        d_ = ImageDraw.Draw(ov)
        d_.rounded_rectangle([bx8,185,bx8+bw8,210], radius=18, fill=(*col,50))
        img = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')

        img = glow(img, num, (bx8+bw8//2, 200), F(100,True), col, col, 30, 'mt')
        d   = ImageDraw.Draw(img)
        uw  = tw(d, unit, F(26,True))
        d.text((bx8+(bw8-uw)//2, 308), unit, font=F(26,True), fill=WHT)
        for j,l in enumerate(desc.split('\n')):
            lw_ = tw(d, l, F(22))
            d.text((bx8+(bw8-lw_)//2, 344+j*28), l, font=F(22), fill=GRY)

    # Chart image
    img = paste(img, 'img_timing.png', 80, 455, 900, 520)

    # Right: detailed breakdown table
    img = card(img, 1010, 455, 850, 520, fill=CARD2, border=DGR, brad=16, bw=1, glow_=False)
    d   = ImageDraw.Draw(img)
    d.text((1030, 470), "Phase-by-Phase Breakdown", font=F(30,True), fill=CYN)
    d.line([1030, 508, 1840, 508], fill=(*CYN3,100), width=1)

    rows8 = [
        ("Attack Detection",    "<1.0 s",    "< 1 second",   CYN),
        ("LLM Analysis",        "2.46 s",    "mean / 200",   PRP),
        ("Mitigation Apply",    "0.17 s",    "mean / 200",   ORG),
        ("Auto-Revert Timer",   "60.02 s",   "±0.09 s",      GRN),
        ("Total E2E Response",  "2.88 s",    "detection+LLM+block", CYN),
        ("Wireshark Agreement", "99.96%",    "5 sessions",   GRN),
        ("False Positive Rate", "< 0.01",    "100 hrs test", ORG),
        ("User Comprehension",  "100%",      "AI explanation", PRP),
    ]
    for ri,(phase,val,note,col) in enumerate(rows8):
        ry8 = 520 + ri*58
        fill8 = CARD if ri%2==0 else CARD2
        img  = card(img, 1020, ry8, 830, 50, fill=fill8, border=col, brad=8, bw=1, glow_=False)
        d    = ImageDraw.Draw(img)
        d.ellipse([1035,ry8+18,1050,ry8+33], fill=col)
        d.text((1058,  ry8+8), phase, font=F(23,True), fill=WHT)
        d.text((1530,  ry8+8), val,   font=F(23,True), fill=col)
        d.text((1680,  ry8+8), note,  font=F(20),      fill=GRY)

    return save_s(img, 8)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WIRESHARK VALIDATION
# ═════════════════════════════════════════════════════════════════════════════
def s09():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, 300, 400, 400, CYN3, 40)

    img = header(img, "Wireshark Validation",
                 "Ground-truth verification across 5 independent 10-minute capture sessions")
    img = snum(img, 9)

    # Big agreement number
    img = glow(img, "99.96%", (W//2, 270), F(130,True), GRN, GRN2, 40, 'mm')
    d   = ImageDraw.Draw(img)
    ctxt(d, "Scapy ↔ Wireshark packet count agreement", 348, F(30), LGR)
    img = hline(img, 396, 200, W-200, GRN, 1, 100)

    # Wireshark chart
    img = paste(img, 'img_wireshark.png', 60, 410, 860, 550)

    # Session table
    sessions = [
        ("Session 1", "12,487", "12,483", "99.97%"),
        ("Session 2", "18,923", "18,915", "99.96%"),
        ("Session 3",  "9,456",  "9,452", "99.96%"),
        ("Session 4", "15,234", "15,229", "99.97%"),
        ("Session 5", "20,891", "20,881", "99.95%"),
    ]
    col_x9 = [980, 1200, 1430, 1660]
    col_h9 = ["Session", "Scapy", "Wireshark", "Agreement"]
    d      = ImageDraw.Draw(img)

    img = card(img, 960, 408, 900, 48, fill=(10,30,14), border=GRN, brad=12, bw=2)
    d   = ImageDraw.Draw(img)
    for ci,(hd,cx) in enumerate(zip(col_h9,col_x9)):
        d.text((cx, 420), hd, font=F(24,True), fill=GRN)

    for ri,row in enumerate(sessions):
        ry9 = 470 + ri*82
        img = card(img, 960, ry9, 900, 70, fill=CARD2 if ri%2==0 else CARD, border=DGR, brad=10, bw=1, glow_=False)
        d   = ImageDraw.Draw(img)
        for ci,(cell,cx) in enumerate(zip(row,col_x9)):
            col9 = GRN if ci==3 else (CYN if ci==0 else LGR)
            bold9= (ci==0 or ci==3)
            d.text((cx, ry9+20), cell, font=F(24,True if bold9 else False), fill=col9)

    # Workflow steps
    steps = ["1. Apply Wireshark filter",
             "2. Launch Network Guardian",
             "3. Inject attack via UI",
             "4. Observe: traffic → block → silence → revert",
             "5. Compare Scapy vs Wireshark packet counts"]
    img = card(img, 960, 892, 900, 148, fill=CARD2, border=DGR, brad=14, bw=1, glow_=False)
    d   = ImageDraw.Draw(img)
    d.text((980, 900), "10-Step Validation Workflow", font=F(26,True), fill=GRY)
    for si,step in enumerate(steps):
        d.text((980, 932+si*30), step, font=F(21), fill=LGR)

    return save_s(img, 9)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — COMPARISON
# ═════════════════════════════════════════════════════════════════════════════
def s10():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, W//2, H//2, 700, CYN3, 30)

    img = header(img, "How We Compare",
                 "Network Guardian vs. existing approaches  (Table VI from paper)")
    img = snum(img, 10)

    # Bar chart image
    img = paste(img, 'img_comparison.png', 60, 172, 800, 500)

    # Comparison table
    columns = ["Approach", "Detect", "Explain", "Auto Block", "Auto Revert", "< 3s Response"]
    rows10  = [
        ("Snort [11]",         "Y", "N",  "N",  "N",  "N"),
        ("Anomaly IDS [12]",   "Y", "N",  "N",  "N",  "N"),
        ("DNN / ML IDS [13]",  "Y", "P",  "N",  "N",  "N"),
        ("LIME / SHAP [14]",   "N", "P",  "N",  "N",  "N"),
        ("Network Guardian",   "Y", "Y",  "Y",  "Y",  "Y"),
    ]
    # Total width must fit: start 860 + sum(col_w10) <= 1900
    col_w10 = [268, 118, 118, 134, 142, 148]   # sum=928, right edge=1788
    col_x10 = [860]
    for cw_ in col_w10[:-1]: col_x10.append(col_x10[-1]+cw_)
    rh10    = 90
    sy10    = 185

    # Header row
    img = card(img, 890, sy10, sum(col_w10), rh10-12, fill=(10,28,60), border=CYN, brad=14, bw=2)
    d   = ImageDraw.Draw(img)
    for ci,(hd,cx,cw_) in enumerate(zip(columns,col_x10,col_w10)):
        hw = tw(d, hd, F(24,True))
        xpos = cx + (cw_ - hw)//2 if ci > 0 else cx+10
        d.text((xpos, sy10+24), hd, font=F(24,True), fill=CYN)

    for ri,row in enumerate(rows10):
        is_ours = (ri == len(rows10)-1)
        ry10 = sy10 + (ri+1)*(rh10-2)
        border10 = GRN if is_ours else DGR
        fill10   = (5,28,12) if is_ours else (CARD2 if ri%2==0 else CARD)
        img = card(img, 890, ry10, sum(col_w10), rh10-10, fill=fill10,
                   border=border10, brad=12, bw=3 if is_ours else 1, glow_=is_ours)
        d   = ImageDraw.Draw(img)
        for ci,(cell,cx,cw_) in enumerate(zip(row,col_x10,col_w10)):
            mid_x = cx + cw_//2
            mid_y = ry10 + (rh10-10)//2
            if ci == 0:
                col10 = GRN if is_ours else LGR
                d.text((cx+10, ry10+20), row[0], font=F(25, is_ours), fill=col10)
            elif cell == "Y":
                draw_check(d, mid_x, mid_y, GRN, 20)
            elif cell == "N":
                draw_xmark(d, mid_x, mid_y, RED, 18)
            else:  # Partial
                pw = tw(d, "Partial", F(21))
                d.text((mid_x-pw//2, mid_y-12), "Partial", font=F(21,True), fill=ORG)

    # Bottom highlight
    bot_y = sy10 + rh10*6 - 6
    img = card(img, 890, bot_y, sum(col_w10), 72, fill=(3,22,10), border=GRN, brad=14, bw=2)
    d   = ImageDraw.Draw(img)
    msg = "Network Guardian: ONLY solution with all 5 capabilities in a single lightweight Python pipeline"
    mw  = tw(d, msg, F(23,True))
    d.text((890 + (sum(col_w10)-mw)//2, bot_y+22), msg, font=F(23,True), fill=GRN)

    # Side stats below chart
    img = paste(img, 'img_attacks3d.png', 60, 682, 800, 350)

    return save_s(img, 10)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CHALLENGES
# ═════════════════════════════════════════════════════════════════════════════
def s11():
    img = make_bg(BG1, (25,15,5))
    img = add_dots(img, 55, (80,50,10), 16)
    img = add_orb(img, W//2, H//2, 650, (50,25,5), 35)

    img = header(img, "Challenges & Limitations",
                 "Honest assessment — understanding boundaries of this research")
    img = snum(img, 11)

    cards11 = [
        (ORG, "T", "Technical Limitations",
         [("CPU-only inference", "2.46s analysis (GPU -- <0.5s)"),
          ("Single-host scope",  "No distributed / multi-network support"),
          ("IPv6 not supported", "IPv4 only in current implementation"),
          ("5 attack types",     "DNS amplification, slowloris not covered"),
          ("Loopback tested",    "Real-network validation is future work")]),
        (RED, "E", "Ethical Considerations",
         [("False positive risk",  "Legitimate users could be temporarily blocked"),
          ("Auto-revert mitigates","60s self-healing limits collateral damage"),
          ("AI explainability",    "LLM outputs are advisory, not definitive"),
          ("Human oversight",      "Critical decisions need human validation"),
          ("Adversarial LLM",      "Prompt injection risk being studied")]),
        (PRP, "F", "Future Research",
         [("GPU acceleration",    "Llama-3 / Mistral on GPU for <0.5s"),
          ("Adaptive thresholds", "EWMA-based dynamic threshold learning"),
          ("SIEM integration",    "Elasticsearch / Splunk / syslog export"),
          ("Threat intelligence", "VirusTotal + MITRE ATT&CK feeds"),
          ("Multi-host",          "Distributed agent + central controller")]),
    ]
    cw11, ch11 = 545, 720
    gap11 = 40
    sx11  = (W-(cw11*3+gap11*2))//2

    for i,(col,badge_l,title,items) in enumerate(cards11):
        bx11 = sx11 + i*(cw11+gap11)
        img  = card(img, bx11, 188, cw11, ch11, fill=CARD2, border=col, brad=20, bw=2)

        # Accent bar
        ov_ = Image.new('RGBA',(W,H),(0,0,0,0))
        d_  = ImageDraw.Draw(ov_)
        d_.rounded_rectangle([bx11,188,bx11+cw11,225], radius=20, fill=(*col,40))
        img  = Image.alpha_composite(img.convert('RGBA'),ov_).convert('RGB')

        img = draw_circle_badge(img, bx11+40, 210, 22, col, badge_l, 20)
        d = ImageDraw.Draw(img)
        d.text((bx11+74, 200), title, font=F(28,True), fill=col)
        d.line([bx11+20, 248, bx11+cw11-20, 248], fill=(*col,80), width=1)

        for j,(lbl,desc) in enumerate(items):
            iy = 264 + j*96
            img = card(img, bx11+16, iy, cw11-32, 82, fill=CARD, border=col, brad=10, bw=1, glow_=False)
            d   = ImageDraw.Draw(img)
            d.text((bx11+32, iy+8),  lbl,  font=F(23,True), fill=col)
            txtw(d, desc, bx11+32, iy+38, F(21), LGR, cw11-64)

    return save_s(img, 11)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ═════════════════════════════════════════════════════════════════════════════
def s12():
    img = make_bg()
    img = add_dots(img)
    img = add_orb(img, W//2, 260, 550, CYN3, 50)
    img = add_orb(img, 200, H-200, 300, GRN2, 40)

    img = header(img, "Conclusion",
                 "Network Guardian demonstrates that local LLM integration into security pipelines is practical and effective")
    img = snum(img, 12)

    # Key achievements
    achievements = [
        (CYN, "<1 SEC",  "Attack Detection",     "Sub-second sliding-window detection"),
        (PRP, "2.46 SEC","AI Explanation",        "Plain English via local TinyLlama"),
        (ORG, "0.17 SEC","Route Block Applied",   "OS-level null route, no shell injection"),
        (GRN, "99%",     "Blocking Success",      "198 of 200 attacks successfully blocked"),
        (CYN, "99.96%",  "Wireshark Agreement",   "Ground-truth packet count validation"),
        (PRP, "100%",    "User Comprehension",     "vs. <30% for raw technical IDS alerts"),
    ]
    # 2-col grid of achievement cards; right panel gets W-1280-60=580px
    aw = 560; ah = 106; gap_a = 16; sx_a = 80
    for i,(col,val,title,desc) in enumerate(achievements):
        row = i//2; col_i = i%2
        bx_a = sx_a + col_i*(aw+gap_a)
        by_a = 192 + row*(ah+gap_a)
        img  = card(img, bx_a, by_a, aw, ah, fill=CARD2, border=col, brad=14, bw=2)

        # Left color stripe
        ov_ = Image.new('RGBA',(W,H),(0,0,0,0))
        d_  = ImageDraw.Draw(ov_)
        d_.rounded_rectangle([bx_a,by_a,bx_a+8,by_a+ah], radius=7, fill=(*col,200))
        img  = Image.alpha_composite(img.convert('RGBA'),ov_).convert('RGB')

        d = ImageDraw.Draw(img)
        d.text((bx_a+24,  by_a+10), val,   font=F(32,True), fill=col)
        d.text((bx_a+152, by_a+10), title, font=F(26,True), fill=WHT)
        d.text((bx_a+152, by_a+46), desc,  font=F(21),      fill=GRY)

    # Right: impact statement — starts at 80+560*2+16*3=1248, width=612
    rp12_x = sx_a + aw*2 + gap_a*3
    rp12_w = W - rp12_x - 60
    img = card(img, rp12_x, 192, rp12_w, 570, fill=CARD2, border=CYN, brad=18, bw=2)
    rx12 = rp12_x + 20
    d    = ImageDraw.Draw(img)
    img  = glow(img, "IMPACT", (rx12+340, 220), F(58,True), CYN, CYN, 22, 'mm')
    d    = ImageDraw.Draw(img)
    impact = [
        "First system to combine all of:",
        "  • Sub-second sliding-window detection",
        "  • Local LLM plain-English explanation",
        "  • Autonomous OS route mitigation",
        "  • 60-second auto-revert daemon",
        "  • in a single Python pipeline",
        "",
        "Especially valuable as an educational",
        "tool — students launch attacks and see",
        "detect >> explain >> block in real-time.",
    ]
    iy12 = 256
    for line_ in impact:
        bold_ = line_.startswith("First") or line_.startswith("Especially")
        d.text((rx12+10, iy12), line_, font=F(24,True if bold_ else False),
               fill=CYN if bold_ else LGR)
        iy12 += 34

    # Research contribution banner
    img = card(img, 80, 728, W-160, 100, fill=(5,28,12), border=GRN, brad=16, bw=3)
    d   = ImageDraw.Draw(img)
    ctxt(d, "Research Contribution: Proved that a lightweight, CPU-only local LLM can be integrated into a", 742, F(26), LGR)
    ctxt(d, "live network security pipeline achieving < 3s end-to-end detection-explanation-mitigation.", 778, F(26,True), GRN)

    # Future one-liner
    img = card(img, 80, 848, W-160, 88, fill=CARD2, border=DGR, brad=14, bw=1, glow_=False)
    d   = ImageDraw.Draw(img)
    ctxt(d, "Next: GPU acceleration (Llama-3) · Adaptive thresholds · SIEM integration · Multi-host deployment", 880, F(26), GRY)

    return save_s(img, 12)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — FUTURE ROADMAP
# ═════════════════════════════════════════════════════════════════════════════
def s13():
    img = make_bg()
    img = add_grid(img, 100, CYN3, 12)
    img = add_orb(img, W//2, H//2, 700, PRP2, 30)

    img = header(img, "Future Roadmap",
                 "Planned extensions to production-ready enterprise deployment")
    img = snum(img, 13)

    roadmap = [
        (CYN,  "Q3 2026", "GPU Acceleration",
         "Replace TinyLlama with Llama-3 8B on GPU.\n"
         "Target: < 0.5s analysis latency.\n"
         "VLLM batching for concurrent threats."),
        (PRP,  "Q3 2026", "Adaptive Thresholds",
         "EWMA-based dynamic threshold learning.\n"
         "Bayesian traffic baseline recalibration.\n"
         "Reduces false positives further."),
        (GRN,  "Q4 2026", "Multi-Host Deployment",
         "Distributed agents per network node.\n"
         "Centralized threat dashboard.\n"
         "Cross-host correlation engine."),
        (ORG,  "Q4 2026", "SIEM Integration",
         "Push ThreatEvents to Elasticsearch.\n"
         "Splunk connector + syslog support.\n"
         "PostgreSQL threat history storage."),
        (CYN,  "2027",    "Threat Intelligence",
         "VirusTotal API for IP reputation.\n"
         "MITRE ATT&CK technique mapping.\n"
         "Automated CVE cross-reference."),
    ]
    # ch13 sized to fit content: header(42) + title(36) + desc(~140) + padding = ~260 + safety margin
    cw13 = 330; ch13 = 420
    gap13 = 24
    sx13  = (W-(cw13*5+gap13*4))//2

    for i,(col,when,title,desc) in enumerate(roadmap):
        bx13 = sx13 + i*(cw13+gap13)
        img  = card(img, bx13, 192, cw13, ch13, fill=CARD2, border=col, brad=18, bw=2)

        ov_ = Image.new('RGBA',(W,H),(0,0,0,0))
        d_  = ImageDraw.Draw(ov_)
        d_.rounded_rectangle([bx13,192,bx13+cw13,230], radius=18, fill=(*col,50))
        img  = Image.alpha_composite(img.convert('RGBA'),ov_).convert('RGB')

        d = ImageDraw.Draw(img)
        ww = tw(d, when, F(22,True))
        d.text((bx13+(cw13-ww)//2, 198), when, font=F(22,True), fill=col)
        d.line([bx13+20, 228, bx13+cw13-20, 228], fill=(*col,80), width=1)
        tw_ = tw(d, title, F(26,True))
        d.text((bx13+(cw13-tw_)//2, 240), title, font=F(26,True), fill=WHT)

        txtw(d, desc, bx13+20, 288, F(22), LGR, cw13-40)

    # Timeline bar at bottom  (just below cards: 192+420+20=632)
    tl_y = 192 + ch13 + 22
    img = hline(img, tl_y, sx13, sx13+(cw13*5+gap13*4), CYN, 3, 120)
    ov2 = Image.new('RGBA',(W,H),(0,0,0,0))
    d2  = ImageDraw.Draw(ov2)
    for i in range(len(roadmap)):
        px = sx13 + i*(cw13+gap13) + cw13//2
        d2.ellipse([px-10, tl_y-10, px+10, tl_y+10], fill=(*roadmap[i][0],220))
    img = Image.alpha_composite(img.convert('RGBA'),ov2).convert('RGB')

    # Bottom banner sits below timeline dot (tl_y + 30)
    bot13 = tl_y + 35
    img = card(img, 80, bot13, W-160, 80, fill=CARD2, border=DGR, brad=14, bw=1, glow_=False)
    d   = ImageDraw.Draw(img)
    ctxt(d, "Network Guardian lays the foundation for a new class of AI-native, self-defending network infrastructure", bot13+24, F(26), GRY)

    return save_s(img, 13)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Q&A
# ═════════════════════════════════════════════════════════════════════════════
def s14():
    img = make_bg(BG1, (6,14,38))
    img = add_dots(img, 50, CYN3, 20)

    # Orbs
    img = add_orb(img, W//2, H//2, 700, CYN3, 55)
    img = add_orb(img, 200, 200, 350, PRP2, 35)
    img = add_orb(img, W-200, H-200, 350, GRN2, 30)

    # Network nodes
    np.random.seed(99)
    nodes14 = [(np.random.randint(80,W-80), np.random.randint(80,H-80)) for _ in range(18)]
    img = add_network(img, nodes14, CYN3, CYN, 22, 140)

    # Top/bottom bars
    ov = Image.new('RGBA',(W,H),(0,0,0,0))
    d_ = ImageDraw.Draw(ov)
    d_.rectangle([0,0,W,5], fill=(*CYN,100))
    d_.rectangle([0,H-5,W,H], fill=(*CYN,100))
    img = Image.alpha_composite(img.convert('RGBA'),ov).convert('RGB')

    # Big question mark
    img = glow(img, "?", (W//2, 280), F(380,True), CYN, CYN, 60, 'mm')
    img = hline(img, 580, 300, W-300, CYN, 2, 120)

    d = ImageDraw.Draw(img)
    ctxt(d, "Thank You!", 608, F(72,True), WHT)
    ctxt(d, "Questions & Discussion", 692, F(38), LGR)

    img = hline(img, 750, 400, W-400, DGR, 1, 100)

    d = ImageDraw.Draw(img)
    ctxt(d, "Arya Y P  ·  Adi Narayan Prasad G  ·  Ashrith S Jain  ·  Abhishek", 768, F(28,True), CYN)
    ctxt(d, "Guide: Prof. Nalina V   |   Dept. of ISE   |   BMSCE, Bengaluru", 810, F(26), GRY)
    ctxt(d, "ICWITE 2026", 852, F(26,True), PRP)

    # 3 quick recap badges at bottom
    badges14 = [("<1s Detection", CYN), ("Local LLM AI", PRP), ("99% Blocking", GRN)]
    bw14 = 280; gap14 = 40
    sx14 = (W-(bw14*3+gap14*2))//2
    for i,(lbl,col) in enumerate(badges14):
        bx14 = sx14 + i*(bw14+gap14)
        img  = card(img, bx14, 908, bw14, 62, fill=CARD2, border=col, brad=31, bw=2)
        d    = ImageDraw.Draw(img)
        lw14 = tw(d, lbl, F(26,True))
        d.text((bx14+(bw14-lw14)//2, 923), lbl, font=F(26,True), fill=col)

    return save_s(img, 14)

# ═════════════════════════════════════════════════════════════════════════════
# RUN ALL SLIDES
# ═════════════════════════════════════════════════════════════════════════════
print("Generating slides...")
PATHS = [
    s01(), s02(), s03(), s04(), s05(), s06(), s07(),
    s08(), s09(), s10(), s11(), s12(), s13(), s14(),
]

# ═════════════════════════════════════════════════════════════════════════════
# BUNDLE INTO PPTX
# ═════════════════════════════════════════════════════════════════════════════
print("\nBundling into PPTX...")
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
blank = prs.slide_layouts[6]

for path in PATHS:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(path, 0, 0, prs.slide_width, prs.slide_height)

prs.save(OUT)
sz = os.path.getsize(OUT) / (1024*1024)
print(f"\nDone! {OUT}")
print(f"Slides: {len(prs.slides)}  |  Size: {sz:.1f} MB")

